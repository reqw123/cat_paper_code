"""
壓縮小工具指令 

pyinstaller -F -w ^
--icon=your_icon.ico ^
--hidden-import=openpyxl ^
--hidden-import=et_xmlfile ^
--hidden-import=fitz ^
--hidden-import=pymupdf ^
--hidden-import=pptx ^
--collect-all=openpyxl ^
--collect-all=PyMuPDF ^
--collect-all=python-pptx ^
--collect-binaries=fitz ^
your_script.py
"""

import tkinter as tk
from tkinter import filedialog, messagebox, simpledialog, ttk
import os
import threading
import queue
import time
import uuid

try:
    from tkinterdnd2 import TkinterDnD, DND_FILES
    DND_AVAILABLE = True
except Exception:
    DND_AVAILABLE = False

try:
    from PIL import Image, ImageTk, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    import cv2
    CV2_AVAILABLE = True
except ImportError:
    CV2_AVAILABLE = False

try:
    import pygame
    pygame.mixer.init()
    PYGAME_AVAILABLE = True
except Exception:
    PYGAME_AVAILABLE = False

try:
    import vlc
    VLC_AVAILABLE = True
except Exception:
    VLC_AVAILABLE = False

try:
    import docx as _docx
    DOCX_AVAILABLE = True
except Exception:
    DOCX_AVAILABLE = False

try:
    import fitz as _fitz   # PyMuPDF
    FITZ_AVAILABLE = True
except Exception:
    FITZ_AVAILABLE = False

try:
    from pptx import Presentation as _Presentation
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

try:
    import openpyxl as _openpyxl
    OPENPYXL_AVAILABLE = True
except Exception:
    OPENPYXL_AVAILABLE = False

try:
    import msoffcrypto
    MSOFFCRYPTO_AVAILABLE = True
except Exception:
    MSOFFCRYPTO_AVAILABLE = False

import json as _json

SUPPORTED_EXTS = {'.jpg', '.jpeg', '.png', '.mp3', '.mp4', '.mov',
                  '.txt', '.json', '.docx', '.pdf', '.pptx',
                  '.xlsx', '.xlsm', '.xls'}
IMAGE_EXTS     = {'.jpg', '.jpeg', '.png'}
VIDEO_EXTS     = {'.mp4', '.mov'}
AUDIO_EXTS     = {'.mp3'}
TEXT_EXTS      = {'.txt', '.json', '.docx', '.pptx'}   # PDF 獨立處理
PDF_EXTS       = {'.pdf'}
EXCEL_EXTS     = {'.xlsx', '.xlsm', '.xls'}

PDF_ZOOM_STEPS = [0.5, 0.75, 1.0, 1.25, 1.5, 2.0, 2.5, 3.0, 4.0]
SEARCH_EXTS    = {'.txt', '.json', '.csv', '.yaml', '.yml', '.md', '.log',
                  '.docx', '.pdf', '.xlsx', '.xlsm', '.xls'}
THUMB_W, THUMB_H = 100, 68   # thumbnail pixel dimensions
THUMB_ITEM_H     = THUMB_H + 20

BG       = "#1e1e1e"
FG       = "#d4d4d4"
ENTRY_BG = "#2d2d2d"
BTN_BG   = "#3c3c3c"
ACCENT   = "#0e639c"
FONT     = ("Segoe UI", 10)


class FileRenamer:
    def __init__(self, root):
        self.root = root
        self.root.title("檔案重新命名工具")
        self.root.geometry("960x740")
        self.root.configure(bg=BG)

        self.files          = []
        self.current_index  = 0
        self.folder_path    = ""
        self._photo_ref     = None

        # 媒體播放狀態
        self._playing       = False
        self._paused        = False
        self._stop_event    = threading.Event()
        self._frame_q       = queue.Queue(maxsize=4)
        self._video_thread  = None
        self._after_id      = None   # root.after handle

        # VLC 播放器（延遲初始化）
        self._vlc_instance  = None
        self._vlc_player    = None

        # 文字字體大小（動態調整，全局生效）
        self._text_font_size = 10

        # PDF 檢視器狀態
        self._pdf_doc       = None   # 目前開啟的 fitz.Document
        self._pdf_path      = None
        self._pdf_page_num  = 0
        self._pdf_total     = 0
        self._pdf_zoom      = 1.5
        self._pdf_cache     = {}     # (path, page, zoom) → (PhotoImage, w, h)
        self._pdf_loading   = False  # 渲染執行緒是否正在跑
        self._pdf_pending   = False  # 渲染中又有新請求
        self._pdf_render_id = None   # debounce after() handle
        self._pdf_render_gen = 0     # 世代號：避免過時渲染結果蓋掉新畫面
        self._pdf_offset_x  = 0     # 拖曳平移量
        self._pdf_offset_y  = 0
        self._pdf_drag_xy   = None  # 拖曳起始座標

        # 縮圖側欄狀態
        self._thumb_photos  = {}            # filepath → PhotoImage（最多保留 200 張）
        self._thumb_slots   = []            # [{bg, ph, lbl, img}, ...]
        self._thumb_gen     = 0             # 世代號：新清單時讓舊執行緒結果作廢
        self._thumb_stop    = threading.Event()
        self._thumb_thread  = None

        # 開啟中的大文字預覽 widgets（Text instances）
        self._open_text_viewers = []
        # per-Text-widget search index for cycling matches
        self._viewer_search_index = {}
        # 每個文字視窗個別的字體大小：Text widget -> size
        self._viewer_font_sizes = {}
        # 最後被點擊（或聚焦）的 Text widget，鍵盤控制會作用在此視窗
        self._last_active_text_widget = None

        # 路徑歷史
        self._path_history = []

        # PDF 密碼快取：filepath → password
        self._pdf_passwords = {}

        self._build_ui()
        self._load_path_history()
        self._bind_keys()
        # 確保初始字體設定套用到內建文字預覽
        self._apply_text_font()
        # 初始狀態：尚未載入文字檔，A+/A- 停用
        self._text_inc_btn.config(state=tk.DISABLED)
        self._text_dec_btn.config(state=tk.DISABLED)
        self.root.protocol("WM_DELETE_WINDOW", self._on_close)
        self._setup_dnd()

    # ── UI 建立 ────────────────────────────────────────────────────────────

    def _build_ui(self):
        # ── 頂部：路徑輸入 ──
        top = tk.Frame(self.root, bg=BG, pady=8)
        top.pack(fill=tk.X, padx=12)

        tk.Label(top, text="路徑:", bg=BG, fg=FG, font=FONT).pack(side=tk.LEFT)
        self.path_var = tk.StringVar()
        self.path_combo = ttk.Combobox(top, textvariable=self.path_var, values=self._path_history,
                                        font=FONT, width=50, state="normal")
        self.path_combo.pack(side=tk.LEFT, padx=6, ipady=4)
        self.path_combo.bind("<Return>", lambda _: self._load_files())
        self.path_combo.bind("<<ComboboxSelected>>", lambda _: self._load_files())

        tk.Button(top, text="資料夾", command=self._browse, bg=BTN_BG, fg=FG,
                  relief=tk.FLAT, font=FONT, padx=8, cursor="hand2").pack(side=tk.LEFT)
        tk.Button(top, text="單一檔案", command=self._browse_file, bg=BTN_BG, fg=FG,
                  relief=tk.FLAT, font=FONT, padx=8, cursor="hand2").pack(side=tk.LEFT, padx=4)
        tk.Button(top, text="載入", command=self._load_files, bg=ACCENT, fg="white",
                  relief=tk.FLAT, font=FONT, padx=8, cursor="hand2").pack(side=tk.LEFT, padx=(0, 6))
        tk.Button(top, text="批次重命名", command=self._open_batch_rename,
                  bg="#7c4a7c", fg="white", relief=tk.FLAT, font=FONT,
                  padx=8, cursor="hand2").pack(side=tk.LEFT)
        tk.Button(top, text="🔍 搜尋", command=self._open_search,
                  bg="#2a4a2a", fg="#9cdcfe", relief=tk.FLAT, font=FONT,
                  padx=8, cursor="hand2").pack(side=tk.LEFT, padx=(4, 0))

        self._fs_btn = tk.Button(top, text="⊡ 全螢幕 [F11]",
                                  command=self._toggle_fullscreen,
                                  bg="#3a2a4a", fg="#c586c0",
                                  relief=tk.FLAT, font=FONT,
                                  padx=8, cursor="hand2")
        self._fs_btn.pack(side=tk.RIGHT, padx=(4, 0))

        # ── 中部：主區域（左縮圖欄 + 右預覽容器）──
        self._main_area = tk.Frame(self.root, bg=BG)
        self._main_area.pack(fill=tk.BOTH, expand=True, padx=12, pady=4)

        # 左側縮圖欄
        self._sidebar_frame = tk.Frame(self._main_area, bg="#1a1a1a", width=120)
        self._sidebar_frame.pack(side=tk.LEFT, fill=tk.Y)
        self._sidebar_frame.pack_propagate(False)

        _sb = tk.Scrollbar(self._sidebar_frame, orient=tk.VERTICAL,
                           bg="#252526", troughcolor="#1a1a1a")
        _sb.pack(side=tk.RIGHT, fill=tk.Y)
        self._thumb_canvas = tk.Canvas(self._sidebar_frame, bg="#1a1a1a",
                                       highlightthickness=0, yscrollcommand=_sb.set)
        self._thumb_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        _sb.config(command=self._thumb_canvas.yview)
        self._thumb_canvas.bind("<Enter>",
            lambda e: self._thumb_canvas.bind_all("<MouseWheel>", self._thumb_on_wheel))
        self._thumb_canvas.bind("<Leave>",
            lambda e: self._thumb_canvas.unbind_all("<MouseWheel>"))

        # 右側共用容器（canvas / 文字 / VLC 視圖互換）
        self._view_container = tk.Frame(self._main_area, bg="#111111")
        self._view_container.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # canvas（圖片 / 影片）
        self.canvas = tk.Canvas(self._view_container, bg="#111111", highlightthickness=0)
        self.canvas.pack(fill=tk.BOTH, expand=True)
        self.canvas_img_id  = self.canvas.create_image(0, 0, anchor=tk.NW)
        self.canvas_text_id = self.canvas.create_text(
            480, 280, text="", fill=FG, font=("Segoe UI", 18), justify=tk.CENTER)

        # 文字預覽框（txt / json / docx / pdf / pptx）
        self._text_frame = tk.Frame(self._view_container, bg="#1e1e1e")
        _tsb = tk.Scrollbar(self._text_frame, bg=BTN_BG, troughcolor=BG)
        _tsb.pack(side=tk.RIGHT, fill=tk.Y)
        self._text_widget = tk.Text(
            self._text_frame, bg="#1e1e1e", fg=FG,
            font=("Consolas", 10), wrap=tk.WORD,
            height=1,           # 關鍵：固定 reqheight = 1行高度
                                # 不論字體多大，Text 向 pack 請求的最小高度
                                # 永遠只有1行，不會把下方的列推出視窗外
                                # expand=True 讓它實際填滿所有剩餘空間
            state=tk.DISABLED, relief=tk.FLAT,
            padx=12, pady=8, yscrollcommand=_tsb.set,
            insertbackground=FG, selectbackground="#264f78")
        self._text_widget.pack(fill=tk.BOTH, expand=True)
        _tsb.config(command=self._text_widget.yview)
        # 初始化內嵌文字視窗的個別字體大小與點擊追蹤
        self._viewer_font_sizes[self._text_widget] = self._text_font_size
        self._last_active_text_widget = self._text_widget
        self._text_widget.bind("<Button-1>", lambda e: self._set_active_text_widget(self._text_widget))

        # VLC 專用 Frame（獨立 HWND，與 canvas 完全隔離）
        self._vlc_frame = tk.Frame(self._view_container, bg="black")
        # Excel 格子視圖
        self._excel_frame = tk.Frame(self._view_container, bg="#1e1e1e")
        # 視圖預設只顯示 canvas，其餘隱藏

        # ── 媒體控制列 (MP3 / 影片時顯示) ──
        self.media_bar = tk.Frame(self.root, bg="#252526", pady=5)
        self.media_bar.pack(fill=tk.X, padx=12)

        self.play_btn = tk.Button(self.media_bar, text="▶  播放", width=9,
                                  command=self._play, bg="#4a7c59", fg="white",
                                  relief=tk.FLAT, font=FONT, cursor="hand2")
        self.play_btn.pack(side=tk.LEFT, padx=4)

        self.pause_btn = tk.Button(self.media_bar, text="⏸  暫停", width=9,
                                   command=self._pause, bg=BTN_BG, fg=FG,
                                   relief=tk.FLAT, font=FONT, cursor="hand2",
                                   state=tk.DISABLED)
        self.pause_btn.pack(side=tk.LEFT, padx=4)

        self.stop_btn = tk.Button(self.media_bar, text="⏹  停止", width=9,
                                  command=self._stop, bg=BTN_BG, fg=FG,
                                  relief=tk.FLAT, font=FONT, cursor="hand2",
                                  state=tk.DISABLED)
        self.stop_btn.pack(side=tk.LEFT, padx=4)

        self.media_status_var = tk.StringVar(value="")
        tk.Label(self.media_bar, textvariable=self.media_status_var,
                 bg="#252526", fg="#9cdcfe", font=("Segoe UI", 9)).pack(side=tk.LEFT, padx=10)

        # 提示：影片播放無聲音
        self.media_note_var = tk.StringVar(value="")
        tk.Label(self.media_bar, textvariable=self.media_note_var,
                 bg="#252526", fg="#ce9178", font=("Segoe UI", 8)).pack(side=tk.RIGHT, padx=8)

        # ── PDF 翻頁 / 縮放控制列 ──
        # 先 pack 後立刻 pack_forget，把位置鎖在 media_bar 之後、info_bar 之前。
        # 之後 _set_pdf_bar() 再 pack/pack_forget 時 Tkinter 會沿用此順序。
        self._pdf_bar = tk.Frame(self.root, bg="#1e2a3a", pady=5)
        self._pdf_bar.pack(fill=tk.X, padx=12)
        self._pdf_bar.pack_forget()

        self._pdf_prev_btn = tk.Button(
            self._pdf_bar, text="◀ 上頁", command=self._pdf_prev_page,
            bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT, padx=8, cursor="hand2")
        self._pdf_prev_btn.pack(side=tk.LEFT, padx=4)

        self._pdf_page_var = tk.StringVar(value="第 0 / 0 頁")
        tk.Label(self._pdf_bar, textvariable=self._pdf_page_var,
                 bg="#1e2a3a", fg="#9cdcfe", font=("Segoe UI", 9), width=14).pack(side=tk.LEFT)

        self._pdf_next_btn = tk.Button(
            self._pdf_bar, text="下頁 ▶", command=self._pdf_next_page,
            bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT, padx=8, cursor="hand2")
        self._pdf_next_btn.pack(side=tk.LEFT, padx=4)

        tk.Frame(self._pdf_bar, bg="#3c4a5a", width=1).pack(side=tk.LEFT, fill=tk.Y, padx=8)

        tk.Button(self._pdf_bar, text="縮小 −", command=self._pdf_zoom_out,
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT, padx=6, cursor="hand2").pack(side=tk.LEFT)
        self._pdf_zoom_var = tk.StringVar(value="150%")
        tk.Label(self._pdf_bar, textvariable=self._pdf_zoom_var,
                 bg="#1e2a3a", fg=FG, font=("Segoe UI", 9), width=5).pack(side=tk.LEFT)
        tk.Button(self._pdf_bar, text="放大 +", command=self._pdf_zoom_in,
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT, padx=6, cursor="hand2").pack(side=tk.LEFT)
        tk.Button(self._pdf_bar, text="適合寬度", command=self._pdf_zoom_fit,
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT, padx=6, cursor="hand2").pack(side=tk.LEFT, padx=(4, 0))

        self._pdf_status_var = tk.StringVar(value="")
        tk.Label(self._pdf_bar, textvariable=self._pdf_status_var,
                 bg="#1e2a3a", fg="#ce9178", font=("Segoe UI", 8)).pack(side=tk.RIGHT, padx=8)

        # ── 文字 / Excel 字體大小列 ──
        # 先 pack 後立刻 pack_forget，把位置鎖在 _pdf_bar 之後、info_bar 之前。
        # pack_propagate(False) + height=36 確保欄高固定，不會因內容字體變動而撐大。
        self._text_bar = tk.Frame(self.root, bg="#1e2a1e", height=36)
        self._text_bar.pack(fill=tk.X, padx=12)
        self._text_bar.pack_forget()
        self._text_bar.pack_propagate(False)

        # place() 讓內容層填滿外框，與外框的 pack 幾何管理器完全隔離。
        _tb_inner = tk.Frame(self._text_bar, bg="#1e2a1e")
        _tb_inner.place(relx=0, rely=0, relwidth=1, relheight=1)

        tk.Label(_tb_inner, text="字體大小:", bg="#1e2a1e", fg=FG,
                 font=("Segoe UI", 9)).grid(row=0, column=0, padx=(10, 4), pady=4, sticky="w")

        # 按鈕不設 width/height（文字單位），改用 padx 控制間距，
        # 避免字體大小改變時 Tk 重算文字單位造成按鈕尺寸改變。
        self._text_dec_btn = tk.Button(
            _tb_inner, text="A−", command=self._text_font_dec,
            bg=BTN_BG, fg="#ce9178", relief=tk.FLAT,
            font=("Segoe UI", 10, "bold"),
            cursor="hand2", padx=10, pady=2)
        self._text_dec_btn.grid(row=0, column=1, padx=2, pady=3)

        self._text_font_var = tk.StringVar(value="10 pt")
        tk.Label(_tb_inner, textvariable=self._text_font_var,
                 bg="#1e2a1e", fg="#9cdcfe", font=("Segoe UI", 9),
                 width=6, anchor="center").grid(row=0, column=2, padx=2, pady=4)

        self._text_inc_btn = tk.Button(
            _tb_inner, text="A+", command=self._text_font_inc,
            bg=BTN_BG, fg="#4ec9b0", relief=tk.FLAT,
            font=("Segoe UI", 10, "bold"),
            cursor="hand2", padx=10, pady=2)
        self._text_inc_btn.grid(row=0, column=3, padx=2, pady=3)

        tk.Label(_tb_inner, text="  Ctrl+／Ctrl−",
                 bg="#1e2a1e", fg="#555555",
                 font=("Segoe UI", 8)).grid(row=0, column=4, padx=4, pady=4, sticky="w")

        # ── 檔案資訊列 ──
        self.info_bar = tk.Frame(self.root, bg="#1a1a1a", pady=4)
        self.info_bar.pack(fill=tk.X, padx=12)

        self.info_var = tk.StringVar(value="請選擇資料夾並按「載入」")
        tk.Label(self.info_bar, textvariable=self.info_var, bg="#1a1a1a", fg="#9cdcfe",
                 font=("Segoe UI", 9)).pack(side=tk.LEFT, padx=8)

        tk.Button(self.info_bar, text="⬡ 開啟原始檔案", command=self._open_with_system,
                  bg="#2a3a2a", fg="#4ec9b0", relief=tk.FLAT, font=("Segoe UI", 9),
                  padx=8, cursor="hand2").pack(side=tk.LEFT, padx=4)

        self.counter_var = tk.StringVar(value="0 / 0")
        tk.Label(self.info_bar, textvariable=self.counter_var, bg="#1a1a1a", fg=FG,
                 font=("Segoe UI", 9)).pack(side=tk.RIGHT, padx=8)

        self._expand_btn = tk.Button(
            self.info_bar, text="大預覽 ↗", command=self._open_text_viewer,
            bg="#5a4a2a", fg="#ffd700", relief=tk.FLAT, font=("Segoe UI", 9),
            padx=8, cursor="hand2")
        # 預設隱藏，文字檔時才顯示

        # ── 底部：重命名 + 導航 ──
        bottom = tk.Frame(self.root, bg=BG, pady=8)
        bottom.pack(fill=tk.X, padx=12)

        tk.Label(bottom, text="新檔名 (不含副檔名):", bg=BG, fg=FG, font=FONT).pack(side=tk.LEFT)
        self.rename_var = tk.StringVar()
        self.rename_entry = tk.Entry(bottom, textvariable=self.rename_var, bg=ENTRY_BG,
                                     fg=FG, insertbackground=FG, relief=tk.FLAT,
                                     font=FONT, width=36)
        self.rename_entry.pack(side=tk.LEFT, padx=6, ipady=4)

        tk.Button(bottom, text="儲存重命名  [Enter]", command=self._rename,
                  bg="#4a7c59", fg="white", relief=tk.FLAT, font=FONT,
                  padx=8, cursor="hand2").pack(side=tk.LEFT)

        tk.Button(bottom, text="→ 下一個 [→]", command=self._next,
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT,
                  padx=8, cursor="hand2").pack(side=tk.RIGHT)
        tk.Button(bottom, text="← 上一個 [←]", command=self._prev,
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT,
                  padx=8, cursor="hand2").pack(side=tk.RIGHT, padx=4)

    def _bind_keys(self):
        self.root.bind("<Left>",   lambda _: self._prev())
        self.root.bind("<Right>",  lambda _: self._next())
        self.root.bind("<Return>", lambda _: self._rename())
        self.root.bind("<space>",  lambda _: self._toggle_pause())
        self.root.bind("<F11>",    lambda _: self._toggle_fullscreen())
        # 全域鍵盤：根據最後被點擊的文字視窗，控制該視窗的字體大小
        self.root.bind_all("<Control-equal>", lambda e: self._global_text_font_inc(e))
        self.root.bind_all("<Control-minus>", lambda e: self._global_text_font_dec(e))

    # ── 資料夾 & 檔案載入 ─────────────────────────────────────────────────

    def _browse(self):
        if not self._ensure_stopped_for_action():
            return
        folder = filedialog.askdirectory(title="選擇資料夾")
        if folder:
            self.path_var.set(folder)
            self._load_files()

    def _browse_file(self):
        if not self._ensure_stopped_for_action():
            return
        exts = " ".join(f"*{e}" for e in sorted(SUPPORTED_EXTS))
        path = filedialog.askopenfilename(
            title="開啟單一檔案",
            filetypes=[
                ("支援的格式", exts),
                ("Excel 活頁簿", "*.xlsx *.xlsm *.xls"),
                ("全部檔案", "*.*"),
            ])
        if path:
            self.path_var.set(path)
            self._load_files()

    # ── 路徑歷史 ────────────────────────────────────────────────────────────

    def _load_path_history(self):
        """從 JSON 檔案載入路徑歷史"""
        history_file = os.path.expanduser("~/.file_renamer_history.json")
        try:
            if os.path.exists(history_file):
                with open(history_file, 'r', encoding='utf-8') as f:
                    data = _json.load(f)
                    self._path_history = data.get('paths', [])[:15]  # 最多 15 筆
            else:
                self._path_history = []
        except Exception:
            self._path_history = []

    def _save_path_history(self):
        """將路徑歷史保存至 JSON 檔案"""
        history_file = os.path.expanduser("~/.file_renamer_history.json")
        try:
            data = {'paths': self._path_history}
            with open(history_file, 'w', encoding='utf-8') as f:
                _json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception:
            pass

    def _add_path_to_history(self, path):
        """新增路徑到歷史，若已存在則移到最前面"""
        if not path or not os.path.isabs(path):
            return
        # 移除已存在的相同路徑
        if path in self._path_history:
            self._path_history.remove(path)
        # 添加到最前面
        self._path_history.insert(0, path)
        # 保持最多 15 筆
        self._path_history = self._path_history[:15]
        # 保存到檔案
        self._save_path_history()
        # 更新 Combobox 的 values
        self.path_combo.config(values=self._path_history)

    def _load_files(self):
        path = self.path_var.get().strip()
        # Ensure any playing media is stopped before loading files
        if not self._ensure_stopped_for_action():
            return

        if os.path.isfile(path):
            # 單一檔案：只加入這個檔案，folder_path 設為父目錄
            if os.path.splitext(path)[1].lower() not in SUPPORTED_EXTS:
                messagebox.showerror("錯誤", f"不支援的檔案類型:\n{os.path.basename(path)}")
                return
            self.folder_path = os.path.dirname(path)
            self.files = [path]

        elif os.path.isdir(path):
            # 資料夾：遞迴掃描所有子目錄
            self.folder_path = path
            found = []
            for dirpath, _, filenames in os.walk(path):
                for fname in filenames:
                    if os.path.splitext(fname)[1].lower() in SUPPORTED_EXTS:
                        found.append(os.path.join(dirpath, fname))
            self.files = sorted(found)

        else:
            messagebox.showerror("錯誤", f"路徑不存在:\n{path}")
            return

        self.current_index = 0
        self._sidebar_rebuild()
        if self.files:
            self._add_path_to_history(path)
            self._show_current()
        else:
            self.info_var.set("未找到支援的檔案 (jpg/png/mp3/mp4/mov/txt/json/docx/pdf/pptx/xlsx)")
            self.counter_var.set("0 / 0")
            self._set_canvas_text("沒有找到支援的檔案")

    # ── 顯示當前檔案 ──────────────────────────────────────────────────────

    def _show_current(self):
        if not self.files:
            return
        self._stop()
        self._thumb_update_highlight()

        filepath = self.files[self.current_index]
        filename = os.path.basename(filepath)
        ext      = os.path.splitext(filename)[1].lower()
        size_kb  = os.path.getsize(filepath) / 1024

        self.info_var.set(f"{filename}    {size_kb:.1f} KB")
        self.counter_var.set(f"{self.current_index + 1} / {len(self.files)}")
        self.rename_var.set(os.path.splitext(filename)[0])

        # A+/A- 只在文字類（含 Excel）檔案時啟用
        is_text_type = ext in TEXT_EXTS or ext in EXCEL_EXTS
        _font_btn_state = tk.NORMAL if is_text_type else tk.DISABLED
        self._text_inc_btn.config(state=_font_btn_state)
        self._text_dec_btn.config(state=_font_btn_state)

        if ext in IMAGE_EXTS:
            self._set_media_bar(visible=False)
            self._set_text_bar(visible=False)
            self._set_pdf_bar(visible=False)
            self._expand_btn.pack_forget()
            self._pdf_cancel_all()
            self._switch_to_canvas_view()
            self._show_image(filepath)
        elif ext in VIDEO_EXTS:
            note = "" if VLC_AVAILABLE else "未安裝 python-vlc，影片播放將靜音"
            self._set_media_bar(visible=True, note=note)
            self._set_text_bar(visible=False)
            self._set_pdf_bar(visible=False)
            self._expand_btn.pack_forget()
            self._pdf_cancel_all()
            self._switch_to_canvas_view()
            self._show_video_thumb(filepath)
        elif ext in AUDIO_EXTS:
            self._set_media_bar(visible=True, note="")
            self._set_text_bar(visible=False)
            self._set_pdf_bar(visible=False)
            self._expand_btn.pack_forget()
            self._pdf_cancel_all()
            self._switch_to_canvas_view()
            self._set_canvas_text(f"♪  音訊檔案\n\n{filename}\n\n{size_kb:.1f} KB")
        elif ext in TEXT_EXTS:
            self._set_media_bar(visible=False)
            self._set_text_bar(visible=True)
            self._set_pdf_bar(visible=False)
            self._expand_btn.config(command=self._open_text_viewer)
            self._expand_btn.pack(side=tk.RIGHT, padx=6)
            self._pdf_cancel_all()
            self._switch_to_text_view()
            self._show_text_preview(filepath)
        elif ext in PDF_EXTS:
            self._set_media_bar(visible=False)
            self._set_text_bar(visible=False)
            self._set_pdf_bar(visible=True)
            self._expand_btn.config(command=self._open_text_viewer)
            self._expand_btn.pack(side=tk.RIGHT, padx=6)
            self._switch_to_canvas_view()
            self._pdf_open(filepath)
        elif ext in EXCEL_EXTS:
            self._set_media_bar(visible=False)
            self._set_pdf_bar(visible=False)
            self._set_text_bar(visible=True)
            self._expand_btn.config(
                command=lambda fp=filepath: self._open_excel_viewer(fp))
            self._expand_btn.pack(side=tk.RIGHT, padx=6)
            self._pdf_cancel_all()
            self._switch_to_excel_view()
            self._show_excel_grid(filepath)

    def _set_media_bar(self, visible, note=""):
        if visible:
            self.media_bar.pack(fill=tk.X, padx=12)
        else:
            self.media_bar.pack_forget()
        self.media_note_var.set(note)
        self.media_status_var.set("⏹ 已停止")
        self.play_btn.config(state=tk.NORMAL)
        self.pause_btn.config(state=tk.DISABLED)
        self.stop_btn.config(state=tk.DISABLED)

    def _ensure_stopped_for_action(self):
        """If media is playing, ask user to stop before performing actions that may block.
        Returns True if it's OK to proceed (either not playing or user agreed and stop completed),
        False if user cancelled the action.
        """
        if not self._playing:
            return True
        res = messagebox.askyesno("播放中", "目前正在播放媒體，是否先停止並繼續操作？")
        if not res:
            return False
        self._stop()
        # 讓 Tk 事件迴圈處理停止後的回呼，最多等待 0.5 秒
        deadline = time.time() + 0.5
        while self._playing and time.time() < deadline:
            self.root.update_idletasks()
            time.sleep(0.02)
        return True

    # ── 預覽 ──────────────────────────────────────────────────────────────

    def _show_image(self, filepath):
        if not PIL_AVAILABLE:
            self._set_canvas_text("需要安裝 Pillow:\npip install Pillow")
            return
        try:
            self.root.update_idletasks()
            cw = max(self.canvas.winfo_width(), 100)
            ch = max(self.canvas.winfo_height(), 100)
            img = Image.open(filepath)
            img.thumbnail((cw, ch), Image.LANCZOS)
            photo = ImageTk.PhotoImage(img)
            self._photo_ref = photo
            cx, cy = cw // 2, ch // 2
            self.canvas.coords(self.canvas_img_id, cx - img.width // 2, cy - img.height // 2)
            self.canvas.itemconfig(self.canvas_img_id, image=photo)
            self.canvas.itemconfig(self.canvas_text_id, text="")
        except Exception as e:
            self._set_canvas_text(f"無法顯示圖片\n{e}")

    def _show_video_thumb(self, filepath):
        if not (CV2_AVAILABLE and PIL_AVAILABLE):
            self._set_canvas_text(f"[影片]\n{os.path.basename(filepath)}\n\n需安裝: pip install opencv-python Pillow")
            return
        try:
            cap = cv2.VideoCapture(filepath)
            ret, frame = cap.read()
            cap.release()
            if not ret:
                self._set_canvas_text(f"[影片]\n{os.path.basename(filepath)}\n(無法讀取縮圖)")
                return
            self._render_cv2_frame(frame)
            self.canvas.itemconfig(self.canvas_text_id, text="▶  點播放鍵開始")
        except Exception as e:
            self._set_canvas_text(f"[影片]\n{e}")

    def _render_cv2_frame(self, frame):
        self.root.update_idletasks()
        cw = max(self.canvas.winfo_width(), 100)
        ch = max(self.canvas.winfo_height(), 100)
        img = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
        img.thumbnail((cw, ch), Image.LANCZOS)
        photo = ImageTk.PhotoImage(img)
        self._photo_ref = photo
        cx, cy = cw // 2, ch // 2
        self.canvas.coords(self.canvas_img_id, cx - img.width // 2, cy - img.height // 2)
        self.canvas.itemconfig(self.canvas_img_id, image=photo)
        self.canvas.itemconfig(self.canvas_text_id, text="")

    def _set_canvas_text(self, text):
        self.canvas.itemconfig(self.canvas_img_id, image="")
        self._photo_ref = None
        self.root.update_idletasks()
        cx = max(self.canvas.winfo_width(), 200) // 2
        cy = max(self.canvas.winfo_height(), 200) // 2
        self.canvas.coords(self.canvas_text_id, cx, cy)
        self.canvas.itemconfig(self.canvas_text_id, text=text)

    # ── PDF 控制列顯示 ────────────────────────────────────────────────────

    def _set_pdf_bar(self, visible):
        if visible:
            self._pdf_bar.pack(fill=tk.X, padx=12)
        else:
            self._pdf_bar.pack_forget()

    def _set_text_bar(self, visible):
        if visible:
            if not self._text_bar.winfo_ismapped():
                self._text_bar.pack(fill=tk.X, padx=12, pady=(0, 2))
        else:
            if self._text_bar.winfo_ismapped():
                self._text_bar.pack_forget()

    # ── 全螢幕切換 ────────────────────────────────────────────────────────────

    def _toggle_fullscreen(self):
        """F11：在最大化（全螢幕）和還原小視窗之間切換。"""
        if self.root.state() == 'zoomed':
            self.root.state('normal')
            self.root.geometry("960x740")
            self._fs_btn.config(text="⊡ 全螢幕 [F11]")
        else:
            self.root.state('zoomed')
            self._fs_btn.config(text="⊟ 還原視窗 [F11]")

    # ── 文字字體大小 ──────────────────────────────────────────────────────────

    def _text_font_inc(self):
        # Increase font size for the currently active text widget.
        self._adjust_font_for_active(1)

    def _text_font_dec(self):
        # Decrease font size for the currently active text widget.
        self._adjust_font_for_active(-1)

    def _global_text_font_inc(self, event=None):
        # 只在文字類型檔案時有效
        if self._text_bar.winfo_ismapped():
            self._adjust_font_for_active(1)

    def _global_text_font_dec(self, event=None):
        # 只在文字類型檔案時有效
        if self._text_bar.winfo_ismapped():
            self._adjust_font_for_active(-1)

    def _set_active_text_widget(self, tw):
        try:
            self._last_active_text_widget = tw
        except Exception:
            pass

    def _adjust_font_for_active(self, delta):
        """Adjust font size for the last active text widget. If it's the embedded
        main text widget, update global text font size and styles; otherwise only
        update that viewer's font."""
        tw = self._last_active_text_widget or self._text_widget
        # Ensure we have a tracked size
        if tw is self._text_widget:
            new = max(7, min(28, self._text_font_size + delta))
            if new != self._text_font_size:
                self._text_font_size = new
                self._apply_text_font()
        else:
            cur = self._viewer_font_sizes.get(tw, self._text_font_size)
            new = max(7, min(28, cur + delta))
            if new != cur:
                self._viewer_font_sizes[tw] = new
                try:
                    tw.config(font=("Consolas", new))
                except Exception:
                    pass

    def _apply_text_font(self):
        sz = self._text_font_size
        self._text_font_var.set(f"{sz} pt")
        # 更新文字 Widget
        self._text_widget.config(font=("Consolas", sz))
        # 更新內嵌文字視窗的追蹤值
        self._viewer_font_sizes[self._text_widget] = sz
        # 更新 Excel Treeview 樣式（rowheight 也跟著動）
        s = ttk.Style()
        s.configure(
            "XL.Treeview",
            font=("Consolas", sz),
            rowheight=24
        )

    # ── PDF 檢視器 ────────────────────────────────────────────────────────

    def _format_docx(self, doc):
        """將 DOCX 內容格式化為易讀格式"""
        lines = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                lines.append("")  # 保留空行
                continue

            # 檢測標題等級
            if para.style and para.style.name.startswith('Heading'):
                try:
                    level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
                except (ValueError, IndexError):
                    level = 1
                prefix = '█' * level + ' '
                lines.append(f"\n{prefix}{text}\n")

            # 檢測列表項目
            elif para.style and para.style.name.startswith('List'):
                indent = '  ' * (para.level or 0)
                lines.append(f"{indent}• {text}")

            # 一般文本
            else:
                lines.append(text)

        # 處理表格（若有）
        if doc.tables:
            for table in doc.tables:
                lines.append("\n┌─ 表格 ─┐")
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append(" │ ".join(cells))
                lines.append("└─────────┘\n")

        return '\n'.join(lines)

    def _format_pptx(self, prs):
        """將 PPTX 內容格式化為易讀格式"""
        slides_out = []

        for slide_num, slide in enumerate(prs.slides, 1):
            # 幻燈片標頭
            slide_header = f"\n{'═' * 60}\n"
            slide_header += f"【投影片 {slide_num}】"
            slide_header += f"\n{'═' * 60}\n"

            slide_content = []

            for shape_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue

                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    # 根據位置推測是否為標題（通常在第一個形狀）
                    level = para.level or 0
                    try:
                        is_title = (shape_idx == 0 and para.font.size and
                                   para.font.size.pt > 24)
                    except (AttributeError, TypeError):
                        is_title = False

                    if is_title:
                        slide_content.append(f"【{text}】")
                    else:
                        indent = '  ' * level
                        slide_content.append(f"{indent}• {text}")

            slide_text = ('\n'.join(slide_content) if slide_content
                         else "  (無文字內容)")
            slides_out.append(slide_header + slide_text)

        return ''.join(slides_out)

    def _try_decrypt_office(self, filepath):
        """
        嘗試解密 Office 檔案（DOCX/XLSX/PPTX）。
        若需密碼，彈出輸入框。返回解密後的 BytesIO，或 None（失敗/取消）
        """
        import io

        if not MSOFFCRYPTO_AVAILABLE:
            return None

        max_attempts = 3
        attempt = 0

        while attempt < max_attempts:
            pwd = simpledialog.askstring(
                "Office 檔案受保護",
                f"檔案「{os.path.basename(filepath)}」已加密。\n"
                f"請輸入密碼（剩餘 {max_attempts - attempt} 次）：",
                show="*"
            )

            if pwd is None:
                return None

            try:
                with open(filepath, 'rb') as f:
                    file = msoffcrypto.OfficeFile(f)
                    file.load_key(password=pwd)

                    decrypted = io.BytesIO()
                    file.decrypt(decrypted)
                    decrypted.seek(0)
                    return decrypted
            except Exception:
                attempt += 1
                if attempt < max_attempts:
                    messagebox.showwarning("密碼錯誤",
                        f"密碼不正確，剩餘 {max_attempts - attempt} 次嘗試")
                else:
                    messagebox.showerror("失敗", "密碼嘗試次數已用完")
                    return None

        return None

    def _ask_pdf_password(self, filepath):
        """彈出密碼輸入對話框，驗證 PDF 密碼。成功返回密碼，失敗或取消返回 None。"""
        max_attempts = 3
        attempt = 0

        while attempt < max_attempts:
            pwd = simpledialog.askstring(
                "PDF 受密碼保護",
                f"檔案「{os.path.basename(filepath)}」受密碼保護。\n"
                f"請輸入密碼（剩餘 {max_attempts - attempt} 次嘗試）：",
                show="*"
            )

            if pwd is None:  # 用戶點擊取消
                return None

            # 驗證密碼
            try:
                doc = _fitz.open(filepath)
                auth_code = doc.authenticate(pwd)
                doc.close()

                if auth_code > 0:  # 密碼正確
                    messagebox.showinfo("成功", "密碼驗證成功")
                    return pwd
                else:  # 密碼錯誤
                    attempt += 1
                    if attempt < max_attempts:
                        messagebox.showwarning("密碼錯誤",
                            f"密碼不正確，剩餘 {max_attempts - attempt} 次嘗試")
                    else:
                        messagebox.showerror("失敗", "密碼嘗試次數已用完")
                        return None
            except Exception as e:
                messagebox.showerror("錯誤", f"驗證失敗：{e}")
                return None

        return None

    def _pdf_open(self, filepath):
        """載入新 PDF：重置狀態、關閉舊文件、清快取（同路徑則保留）。"""
        self._pdf_cancel_all()
        if self._pdf_path != filepath:
            if self._pdf_doc:
                try:
                    self._pdf_doc.close()
                except Exception:
                    pass
            self._pdf_doc   = None
            self._pdf_cache = {}
            self._pdf_path  = filepath
            self._pdf_zoom  = 1.5

        # 檢查 PDF 是否需要密碼
        try:
            doc = _fitz.open(filepath)
            if doc.needs_pass:
                doc.close()
                # 檢查是否已有儲存的密碼
                if filepath in self._pdf_passwords:
                    pwd = self._pdf_passwords[filepath]
                    doc = _fitz.open(filepath)
                    auth_code = doc.authenticate(pwd)
                    doc.close()
                    if auth_code == 0:
                        # 已儲存密碼失效，重新詢問
                        del self._pdf_passwords[filepath]
                        pwd = self._ask_pdf_password(filepath)
                        if pwd is None:
                            self._set_canvas_text("用戶取消密碼輸入")
                            return
                        self._pdf_passwords[filepath] = pwd
                else:
                    # 未曾輸入過，彈出對話框
                    pwd = self._ask_pdf_password(filepath)
                    if pwd is None:
                        self._set_canvas_text("用戶取消密碼輸入")
                        return
                    self._pdf_passwords[filepath] = pwd
            else:
                doc.close()
        except Exception as e:
            self._set_canvas_text(f"PDF 開啟失敗：{e}")
            return

        self._pdf_zoom_var.set(f"{int(self._pdf_zoom * 100)}%")
        self._pdf_page_num = 0
        self._pdf_offset_x = 0
        self._pdf_offset_y = 0
        self._pdf_schedule_render()

    def _pdf_cancel_all(self):
        """取消所有待渲染請求（切換檔案時呼叫）。"""
        if self._pdf_render_id:
            self.root.after_cancel(self._pdf_render_id)
            self._pdf_render_id = None
        self._pdf_pending  = False
        self._pdf_loading  = False   # 世代號遞增後進行中的執行緒結果作廢，旗標必須歸零
        self._pdf_render_gen += 1
        self._pdf_unbind_canvas()

    def _pdf_schedule_render(self):
        """Debounce：快速連按時只保留最後一次請求。"""
        if self._pdf_render_id:
            self.root.after_cancel(self._pdf_render_id)
        self._pdf_render_id = self.root.after(80, self._pdf_start_render)

    def _pdf_start_render(self):
        self._pdf_render_id = None
        if not FITZ_AVAILABLE:
            self._set_canvas_text("需安裝 PyMuPDF:\npip install PyMuPDF")
            return
        if self._pdf_loading:
            self._pdf_pending = True   # 渲染完後再跑一次
            return

        # ── 快取命中 → 直接顯示（無需執行緒）──
        cache_key = (self._pdf_path, self._pdf_page_num, round(self._pdf_zoom, 2))
        if cache_key in self._pdf_cache:
            photo, iw, ih = self._pdf_cache[cache_key]
            self._pdf_apply_render(photo, iw, ih, cache_key, self._pdf_render_gen)
            return

        # ── 背景執行緒渲染 ──
        self._pdf_loading = True
        self._pdf_pending = False
        self._pdf_status_var.set("載入中…")
        # 停用翻頁避免連點
        self._pdf_prev_btn.config(state=tk.DISABLED)
        self._pdf_next_btn.config(state=tk.DISABLED)

        doc_path  = self._pdf_path
        page_num  = self._pdf_page_num
        zoom      = self._pdf_zoom
        gen       = self._pdf_render_gen

        def worker():
            try:
                doc  = _fitz.open(doc_path)

                # 若需密碼，先驗證
                if doc.needs_pass:
                    if doc_path in self._pdf_passwords:
                        pwd = self._pdf_passwords[doc_path]
                        auth_code = doc.authenticate(pwd)
                        if auth_code == 0:
                            raise Exception("PDF 密碼無效")
                    else:
                        raise Exception("PDF 需密碼但未提供")

                total = len(doc)
                page = doc[page_num]
                mat  = _fitz.Matrix(zoom, zoom)
                pix  = page.get_pixmap(matrix=mat, alpha=False)
                img  = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                doc.close()
                # PhotoImage 必須在主執行緒建立
                self.root.after(0, lambda: self._pdf_make_photo(img, pix.width, pix.height,
                                                                 cache_key, gen, total))
            except Exception as e:
                self.root.after(0, lambda: self._pdf_render_error(str(e), gen))

        threading.Thread(target=worker, daemon=True).start()

    def _pdf_make_photo(self, img, iw, ih, cache_key, gen, total):
        """主執行緒把 PIL Image 轉成 PhotoImage（tkinter 要求）再顯示。"""
        if gen != self._pdf_render_gen:
            self._pdf_loading = False
            return
        photo = ImageTk.PhotoImage(img)
        # 快取上限 12 張
        if len(self._pdf_cache) >= 12:
            oldest = next(iter(self._pdf_cache))
            del self._pdf_cache[oldest]
        self._pdf_cache[cache_key] = (photo, iw, ih)
        self._pdf_total = total
        self._pdf_apply_render(photo, iw, ih, cache_key, gen)

    def _pdf_apply_render(self, photo, iw, ih, cache_key, gen):
        """把已渲染圖片放到 canvas，更新 UI 狀態。"""
        if gen != self._pdf_render_gen:
            self._pdf_loading = False
            return
        self.root.update_idletasks()
        cw = max(self.canvas.winfo_width(), 1)
        ch = max(self.canvas.winfo_height(), 1)

        x = cw // 2 - iw // 2 + self._pdf_offset_x
        y = ch // 2 - ih // 2 + self._pdf_offset_y

        self._photo_ref = photo
        self.canvas.coords(self.canvas_img_id, x, y)
        self.canvas.itemconfig(self.canvas_img_id, image=photo)
        self.canvas.itemconfig(self.canvas_text_id, text="")

        pg  = self._pdf_page_num
        tot = self._pdf_total
        self._pdf_page_var.set(f"第 {pg + 1} / {tot} 頁")
        self._pdf_prev_btn.config(state=tk.NORMAL if pg > 0     else tk.DISABLED)
        self._pdf_next_btn.config(state=tk.NORMAL if pg < tot-1 else tk.DISABLED)
        self._pdf_status_var.set("")
        self._pdf_bind_canvas()

        self._pdf_loading = False
        if self._pdf_pending:
            self._pdf_schedule_render()

    def _pdf_render_error(self, msg, gen):
        if gen != self._pdf_render_gen:
            return
        self._pdf_loading = False
        self._pdf_status_var.set(f"錯誤：{msg[:40]}")

    # ── PDF 翻頁 ──────────────────────────────────────────────────────────

    def _pdf_prev_page(self):
        if self._pdf_page_num > 0:
            self._pdf_page_num -= 1
            self._pdf_offset_x = 0
            self._pdf_offset_y = 0
            self._pdf_schedule_render()

    def _pdf_next_page(self):
        if self._pdf_page_num < self._pdf_total - 1:
            self._pdf_page_num += 1
            self._pdf_offset_x = 0
            self._pdf_offset_y = 0
            self._pdf_schedule_render()

    # ── PDF 縮放 ──────────────────────────────────────────────────────────

    def _pdf_zoom_in(self):
        idx = self._pdf_zoom_idx()
        if idx < len(PDF_ZOOM_STEPS) - 1:
            self._pdf_zoom = PDF_ZOOM_STEPS[idx + 1]
            self._pdf_zoom_var.set(f"{int(self._pdf_zoom * 100)}%")
            self._pdf_offset_x = 0
            self._pdf_offset_y = 0
            self._pdf_schedule_render()

    def _pdf_zoom_out(self):
        idx = self._pdf_zoom_idx()
        if idx > 0:
            self._pdf_zoom = PDF_ZOOM_STEPS[idx - 1]
            self._pdf_zoom_var.set(f"{int(self._pdf_zoom * 100)}%")
            self._pdf_offset_x = 0
            self._pdf_offset_y = 0
            self._pdf_schedule_render()

    def _pdf_zoom_fit(self):
        """依 canvas 寬度自動計算縮放比。"""
        if not FITZ_AVAILABLE or not self._pdf_path:
            return
        try:
            doc  = _fitz.open(self._pdf_path)
            rect = doc[self._pdf_page_num].rect
            doc.close()
        except Exception:
            return
        self.root.update_idletasks()
        cw = max(self.canvas.winfo_width(), 100)
        self._pdf_zoom = round(min(max((cw - 24) / rect.width, 0.3), 4.0), 2)
        self._pdf_zoom_var.set(f"{int(self._pdf_zoom * 100)}%")
        self._pdf_offset_x = 0
        self._pdf_offset_y = 0
        self._pdf_schedule_render()

    def _pdf_zoom_idx(self):
        return min(range(len(PDF_ZOOM_STEPS)),
                   key=lambda i: abs(PDF_ZOOM_STEPS[i] - self._pdf_zoom))

    # ── PDF canvas 事件 ───────────────────────────────────────────────────

    def _pdf_bind_canvas(self):
        self.canvas.bind("<Button-1>",       self._pdf_drag_start)
        self.canvas.bind("<B1-Motion>",      self._pdf_drag_move)
        self.canvas.bind("<ButtonRelease-1>",self._pdf_drag_end)
        self.canvas.bind("<MouseWheel>",     self._pdf_wheel)
        self.canvas.config(cursor="fleur")

    def _pdf_unbind_canvas(self):
        for ev in ("<Button-1>", "<B1-Motion>", "<ButtonRelease-1>", "<MouseWheel>"):
            try:
                self.canvas.unbind(ev)
            except Exception:
                pass
        self.canvas.config(cursor="")

    def _pdf_drag_start(self, event):
        self._pdf_drag_xy = (event.x, event.y)

    def _pdf_drag_move(self, event):
        if self._pdf_drag_xy is None:
            return
        dx = event.x - self._pdf_drag_xy[0]
        dy = event.y - self._pdf_drag_xy[1]
        self._pdf_drag_xy = (event.x, event.y)
        self._pdf_offset_x += dx
        self._pdf_offset_y += dy
        self.canvas.move(self.canvas_img_id, dx, dy)

    def _pdf_drag_end(self, _event):
        self._pdf_drag_xy = None

    def _pdf_wheel(self, event):
        if event.delta > 0:
            self._pdf_zoom_in()
        else:
            self._pdf_zoom_out()

    # ── 視圖切換 ──────────────────────────────────────────────────────────

    def _switch_to_canvas_view(self):
        self._text_frame.pack_forget()
        self._vlc_frame.pack_forget()
        self._excel_frame.pack_forget()
        self.canvas.pack(fill=tk.BOTH, expand=True)

    def _switch_to_vlc_view(self):
        self.canvas.pack_forget()
        self._text_frame.pack_forget()
        self._excel_frame.pack_forget()
        self._vlc_frame.pack(fill=tk.BOTH, expand=True)

    def _switch_to_text_view(self):
        self.canvas.pack_forget()
        self._vlc_frame.pack_forget()
        self._excel_frame.pack_forget()
        self._text_frame.pack(fill=tk.BOTH, expand=True)

    def _switch_to_excel_view(self):
        self.canvas.pack_forget()
        self._text_frame.pack_forget()
        self._vlc_frame.pack_forget()
        self._excel_frame.pack(fill=tk.BOTH, expand=True)

    # ── Excel 格子預覽 ────────────────────────────────────────────────────────

    def _show_excel_grid(self, filepath):
        """在 _excel_frame 內建立工作表頁籤 + Treeview 格子。"""
        for w in self._excel_frame.winfo_children():
            w.destroy()

        if not OPENPYXL_AVAILABLE:
            tk.Label(self._excel_frame,
                     text="需安裝 openpyxl:\npip install openpyxl",
                     bg="#1e1e1e", fg="#f44747", font=FONT).pack(pady=30)
            return

        try:
            wb = _openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            sheet_names = wb.sheetnames
            wb.close()
        except Exception as e:
            tk.Label(self._excel_frame, text=f"無法開啟 Excel:\n{e}",
                     bg="#1e1e1e", fg="#f44747", font=FONT).pack(pady=30)
            return

        tab_bar   = tk.Frame(self._excel_frame, bg="#252526")
        tab_bar.pack(fill=tk.X, side=tk.TOP)
        grid_host = tk.Frame(self._excel_frame, bg="#1e1e1e")
        grid_host.pack(fill=tk.BOTH, expand=True, side=tk.TOP)

        self._excel_tab_btns = {}

        def _activate(sname):
            for s, b in self._excel_tab_btns.items():
                b.config(bg="#217346" if s == sname else "#2d2d2d",
                         fg="white"  if s == sname else "#9cdcfe")
            self._excel_load_sheet(filepath, sname, grid_host)

        for sname in sheet_names:
            btn = tk.Button(tab_bar, text=sname, relief=tk.FLAT,
                            bg="#2d2d2d", fg="#9cdcfe",
                            font=("Segoe UI", 9), padx=12, pady=3,
                            cursor="hand2",
                            command=lambda s=sname: _activate(s))
            btn.pack(side=tk.LEFT, padx=1, pady=2)
            self._excel_tab_btns[sname] = btn

        _activate(sheet_names[0])

    def _excel_load_sheet(self, filepath, sheet_name, grid_host):
        """清除 grid_host 並載入指定工作表的 Treeview。"""
        for w in grid_host.winfo_children():
            w.destroy()

        try:
            wb = _openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            ws = wb[sheet_name]
            total_rows = ws.max_row or 0
            total_cols = ws.max_column or 0
            rows_raw = []
            for row in ws.iter_rows(min_row=1, max_row=501, max_col=60, values_only=True):
                rows_raw.append(['' if c is None else str(c) for c in row])
            wb.close()
        except Exception as e:
            tk.Label(grid_host, text=f"讀取失敗: {e}",
                     bg="#1e1e1e", fg="#f44747", font=FONT).pack(pady=20)
            return

        if not rows_raw:
            tk.Label(grid_host, text="(空白工作表)",
                     bg="#1e1e1e", fg=FG, font=FONT).pack(pady=20)
            return

        n_raw_cols = max(len(r) for r in rows_raw)
        used = [ci for ci in range(n_raw_cols)
                if any(ci < len(r) and r[ci] for r in rows_raw)]
        if not used:
            tk.Label(grid_host, text="(全為空白)",
                     bg="#1e1e1e", fg=FG, font=FONT).pack(pady=20)
            return

        rows_trimmed = [[r[ci] if ci < len(r) else '' for ci in used]
                        for r in rows_raw]
        n_cols   = len(used)
        headers  = rows_trimmed[0]
        data_rows = rows_trimmed[1:]

        col_px = []
        for ci in range(n_cols):
            max_len = max(
                len(headers[ci]) if ci < len(headers) else 0,
                max((len(r[ci]) if ci < len(r) else 0)
                    for r in data_rows) if data_rows else 0)
            col_px.append(min(max(max_len * 8 + 14, 64), 220))

        col_ids = [f"c{ci}" for ci in range(n_cols)]

        vsb = tk.Scrollbar(grid_host, orient=tk.VERTICAL,   bg=BTN_BG, troughcolor=BG)
        vsb.pack(side=tk.RIGHT, fill=tk.Y)
        hsb = tk.Scrollbar(grid_host, orient=tk.HORIZONTAL, bg=BTN_BG, troughcolor=BG)
        hsb.pack(side=tk.BOTTOM, fill=tk.X)

        shown = min(len(data_rows), 500)
        note  = (f"  顯示 {shown} / {max(total_rows-1,0)} 列　"
                 f"{n_cols} / {total_cols} 欄"
                 + ("　(截斷前 500 列)" if total_rows > 501 else ""))
        tk.Label(grid_host, text=note, bg="#252526", fg="#9cdcfe",
                 font=("Segoe UI", 8), anchor="w", padx=8
                 ).pack(side=tk.BOTTOM, fill=tk.X)

        style = ttk.Style()
        style.configure("XL.Treeview",
            background="#1e1e1e", foreground=FG,
            fieldbackground="#1e1e1e", rowheight=22,
            font=("Consolas", 9))
        style.configure("XL.Treeview.Heading",
            background="#217346", foreground="white",
            font=("Segoe UI", 9, "bold"), relief="flat")
        style.map("XL.Treeview",
            background=[("selected", "#264f78")],
            foreground=[("selected", "white")])

        tree = ttk.Treeview(grid_host, columns=col_ids, show="headings",
                            style="XL.Treeview",
                            yscrollcommand=vsb.set, xscrollcommand=hsb.set)
        vsb.config(command=tree.yview)
        hsb.config(command=tree.xview)
        tree.pack(fill=tk.BOTH, expand=True)

        for ci, col_id in enumerate(col_ids):
            hdr = (headers[ci] if ci < len(headers) and headers[ci]
                   else (chr(65 + ci) if ci < 26 else f"Col{ci+1}"))
            tree.heading(col_id, text=hdr, anchor="w")
            tree.column(col_id, width=col_px[ci], minwidth=40, anchor="w")

        for ri, row in enumerate(data_rows):
            vals = [row[ci] if ci < len(row) else '' for ci in range(n_cols)]
            tree.insert("", tk.END, values=vals,
                        tags=("even" if ri % 2 == 0 else "odd",))

        tree.tag_configure("odd",  background="#1e1e1e", foreground=FG)
        tree.tag_configure("even", background="#252526", foreground=FG)

    def _open_excel_viewer(self, filepath):
        """大預覽：獨立視窗顯示 Excel Treeview。"""
        if not OPENPYXL_AVAILABLE:
            self._open_text_viewer()
            return
        try:
            wb = _openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            sheet_names = wb.sheetnames
            wb.close()
        except Exception as e:
            messagebox.showerror("錯誤", f"無法開啟 Excel 檔案:\n{e}")
            return

        win = tk.Toplevel(self.root)
        win.title(f"大預覽 — {os.path.basename(filepath)}")
        win.geometry("1200x820")
        win.configure(bg=BG)

        toolbar = tk.Frame(win, bg=BG, pady=6)
        toolbar.pack(fill=tk.X, padx=12)
        tk.Label(toolbar, text=f"工作表: {len(sheet_names)} 個",
                 bg=BG, fg="#9cdcfe", font=("Segoe UI", 9)).pack(side=tk.LEFT)
        tk.Button(toolbar, text="用預設程式開啟",
                  command=lambda: self._open_with_system(filepath),
                  bg="#217346", fg="white", relief=tk.FLAT, font=FONT,
                  padx=10, cursor="hand2").pack(side=tk.RIGHT, padx=6)
        tk.Button(toolbar, text="關閉", command=win.destroy,
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT,
                  padx=10, cursor="hand2").pack(side=tk.RIGHT)

        container = tk.Frame(win, bg=BG)
        container.pack(fill=tk.BOTH, expand=True, padx=12, pady=(0, 8))

        tab_bar   = tk.Frame(container, bg="#252526")
        tab_bar.pack(fill=tk.X, side=tk.TOP)
        grid_host = tk.Frame(container, bg="#1e1e1e")
        grid_host.pack(fill=tk.BOTH, expand=True)

        tab_btns = {}

        def _activate(sname):
            for s, b in tab_btns.items():
                b.config(bg="#217346" if s == sname else "#2d2d2d",
                         fg="white"  if s == sname else "#9cdcfe")
            self._excel_load_sheet(filepath, sname, grid_host)

        for sname in sheet_names:
            btn = tk.Button(tab_bar, text=sname, relief=tk.FLAT,
                            bg="#2d2d2d", fg="#9cdcfe",
                            font=("Segoe UI", 9), padx=12, pady=3,
                            cursor="hand2",
                            command=lambda s=sname: _activate(s))
            btn.pack(side=tk.LEFT, padx=1, pady=2)
            tab_btns[sname] = btn

        _activate(sheet_names[0])

    def _open_with_system(self, filepath=None):
        """用作業系統預設程式開啟檔案（Windows）。"""
        if filepath is None:
            if not self.files:
                return
            filepath = self.files[self.current_index]
        try:
            os.startfile(os.path.normpath(filepath))
        except AttributeError:
            messagebox.showerror("錯誤", "此功能僅支援 Windows")
        except Exception as e:
            messagebox.showerror("錯誤", f"無法開啟檔案:\n{e}")

    # ── 文字檔預覽 ────────────────────────────────────────────────────────

    def _load_docx_content(self, filepath):
        """安全地讀取 DOCX，自動處理加密情況"""
        # 先檢查是否加密
        if MSOFFCRYPTO_AVAILABLE:
            try:
                with open(filepath, 'rb') as f:
                    file = msoffcrypto.OfficeFile(f)
                    if file.is_encrypted():
                        # 檔案已加密，彈出密碼對話框
                        decrypted_io = self._try_decrypt_office(filepath)
                        if decrypted_io is None:
                            return "用戶取消或密碼錯誤"
                        try:
                            doc = _docx.Document(decrypted_io)
                        except Exception as e:
                            return f"解密後仍無法讀取：{e}"
                        return self._format_docx(doc)
            except Exception:
                pass  # 若檢查失敗，嘗試正常打開

        # 嘗試正常打開
        try:
            doc = _docx.Document(filepath)
        except Exception as e:
            # 可能是加密但無法檢測，嘗試解密
            if "需要密碼" in str(e) or "encrypted" in str(e).lower():
                decrypted_io = self._try_decrypt_office(filepath)
                if decrypted_io is None:
                    return "用戶取消或密碼錯誤"
                try:
                    doc = _docx.Document(decrypted_io)
                except Exception as e2:
                    return f"解密後仍無法讀取：{e2}"
            else:
                return f"讀取失敗：{e}"

        return self._format_docx(doc)

    def _load_pptx_content(self, filepath):
        """安全地讀取 PPTX，自動處理加密情況"""
        # 先檢查是否加密
        if MSOFFCRYPTO_AVAILABLE:
            try:
                with open(filepath, 'rb') as f:
                    file = msoffcrypto.OfficeFile(f)
                    if file.is_encrypted():
                        # 檔案已加密，彈出密碼對話框
                        decrypted_io = self._try_decrypt_office(filepath)
                        if decrypted_io is None:
                            return "用戶取消或密碼錯誤"
                        try:
                            prs = _Presentation(decrypted_io)
                        except Exception as e:
                            return f"解密後仍無法讀取：{e}"
                        return self._format_pptx(prs)
            except Exception:
                pass  # 若檢查失敗，嘗試正常打開

        # 嘗試正常打開
        try:
            prs = _Presentation(filepath)
        except Exception as e:
            # 可能是加密但無法檢測，嘗試解密
            if "需要密碼" in str(e) or "encrypted" in str(e).lower():
                decrypted_io = self._try_decrypt_office(filepath)
                if decrypted_io is None:
                    return "用戶取消或密碼錯誤"
                try:
                    prs = _Presentation(decrypted_io)
                except Exception as e2:
                    return f"解密後仍無法讀取：{e2}"
            else:
                return f"讀取失敗：{e}"

        return self._format_pptx(prs)

    def _load_text_content(self, filepath):
        ext = os.path.splitext(filepath)[1].lower()
        try:
            if ext == '.txt':
                with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read()
            elif ext == '.json':
                with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                    raw = _json.load(f)
                return _json.dumps(raw, indent=2, ensure_ascii=False)
            elif ext == '.docx':
                if not DOCX_AVAILABLE:
                    return "需安裝 python-docx:\npip install python-docx"
                return self._load_docx_content(filepath)
            elif ext == '.pdf':
                if not FITZ_AVAILABLE:
                    return "需安裝 PyMuPDF:\npip install PyMuPDF"
                doc = _fitz.open(filepath)
                # 若需密碼，先驗證
                if doc.needs_pass:
                    if filepath in self._pdf_passwords:
                        pwd = self._pdf_passwords[filepath]
                        auth_code = doc.authenticate(pwd)
                        if auth_code == 0:
                            return "PDF 密碼無效"
                    else:
                        return "PDF 受密碼保護，請先在主視窗中預覽"
                pages = [doc[i].get_text() for i in range(len(doc))]
                doc.close()
                return '\n\n'.join(pages)
            elif ext == '.pptx':
                if not PPTX_AVAILABLE:
                    return "需安裝 python-pptx:\npip install python-pptx"
                return self._load_pptx_content(filepath)
            elif ext in EXCEL_EXTS:
                if not OPENPYXL_AVAILABLE:
                    return "需安裝 openpyxl:\npip install openpyxl"
                wb = _openpyxl.load_workbook(filepath, read_only=True, data_only=True)
                sheets_out = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    total_rows = ws.max_row or 0
                    total_cols = ws.max_column or 0

                    # 第一遍：收集所有行（最多 1000 列）
                    rows_raw = []
                    for row in ws.iter_rows(values_only=True):
                        rows_raw.append(['' if c is None else str(c) for c in row])
                        if len(rows_raw) >= 1000:
                            break
                    if not rows_raw:
                        sheets_out.append(
                            f"╔═══ 工作表: {sheet_name} ═══╗\n(空白)")
                        continue

                    # 找出「有內容的欄」索引（跳過全空欄）
                    max_cols = max(len(r) for r in rows_raw)
                    used_cols = [ci for ci in range(max_cols)
                                 if any(ci < len(r) and r[ci] for r in rows_raw)]
                    if not used_cols:
                        sheets_out.append(
                            f"╔═══ 工作表: {sheet_name} ═══╗\n(全為空白)")
                        continue

                    # 計算每個有效欄的顯示寬度（最多 30 字元）
                    col_w = {}
                    for ci in used_cols:
                        col_w[ci] = max(
                            min(len(r[ci]), 30) if ci < len(r) else 0
                            for r in rows_raw)
                        col_w[ci] = max(col_w[ci], 1)

                    # 格式化每行（只輸出 used_cols 欄）
                    lines = []
                    for ri, row in enumerate(rows_raw):
                        cells = []
                        for ci in used_cols:
                            val = (row[ci][:30] if ci < len(row) else '')
                            cells.append(val.ljust(col_w[ci]))
                        lines.append('  ' + ' │ '.join(cells))
                        if ri == 0:    # 標頭後加分隔線
                            sep = '──' + '─┼─'.join(
                                '─' * col_w[ci] for ci in used_cols)
                            lines.append(sep)

                    truncated = f'  … 共 {total_rows} 列，截斷顯示前 1000 列' \
                                if total_rows > 1000 else ''
                    sheets_out.append(
                        f"╔═══ 工作表: {sheet_name}"
                        f"  ({total_rows} 列 × {total_cols} 欄，顯示 {len(used_cols)} 個有效欄) ═══╗\n"
                        + '\n'.join(lines)
                        + (f'\n{truncated}' if truncated else ''))
                wb.close()
                return '\n\n'.join(sheets_out)
        except Exception as e:
            return f"讀取失敗：{e}"
        return ""

    def _show_text_preview(self, filepath):
        content = self._load_text_content(filepath)
        self._current_text_content = content   # 供大預覽使用
        self._current_text_filepath = filepath

        # 只顯示前 300 行作為縮圖預覽
        lines  = content.splitlines()
        preview = '\n'.join(lines[:300])
        if len(lines) > 300:
            preview += f"\n\n… 還有 {len(lines) - 300} 行，按「大預覽 ↗」查看全文"

        tw = self._text_widget
        tw.config(state=tk.NORMAL)
        tw.delete('1.0', tk.END)
        tw.insert(tk.END, preview)
        tw.config(state=tk.DISABLED)
        tw.yview_moveto(0)

    def _open_text_viewer(self):
        if not self.files or self.current_index >= len(self.files):
            return
        filepath = self.files[self.current_index]
        # 使用快取內容（同檔案）；否則重新讀取（含 PDF）
        if (getattr(self, '_current_text_filepath', None) == filepath
                and hasattr(self, '_current_text_content')):
            content = self._current_text_content
        else:
            content = self._load_text_content(filepath)
            if not content:
                content = "(無法讀取文字內容)"

        win = tk.Toplevel(self.root)
        win.title(f"大預覽 — {os.path.basename(filepath)}")
        win.geometry("1000x720")
        win.configure(bg=BG)

        # 工具列
        toolbar = tk.Frame(win, bg=BG, pady=6)
        toolbar.pack(fill=tk.X, padx=12)

        # 搜尋
        tk.Label(toolbar, text="搜尋:", bg=BG, fg=FG, font=FONT).pack(side=tk.LEFT)
        search_var = tk.StringVar()
        search_entry = tk.Entry(toolbar, textvariable=search_var, bg=ENTRY_BG, fg=FG,
                                insertbackground=FG, relief=tk.FLAT, font=FONT, width=28)
        search_entry.pack(side=tk.LEFT, padx=4, ipady=3)

        match_var = tk.StringVar(value="")
        match_lbl = tk.Label(toolbar, textvariable=match_var, bg=BG, fg="#9cdcfe",
                             font=("Segoe UI", 9))
        match_lbl.pack(side=tk.LEFT, padx=6)

        tk.Button(toolbar, text="關閉", command=win.destroy,
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT,
                  padx=10, cursor="hand2").pack(side=tk.RIGHT)

        # 文字區
        frame = tk.Frame(win, bg=BG)
        frame.pack(fill=tk.BOTH, expand=True, padx=12, pady=(0, 8))

        vsb = tk.Scrollbar(frame, bg=BTN_BG, troughcolor=BG)
        vsb.pack(side=tk.RIGHT, fill=tk.Y)
        hsb = tk.Scrollbar(frame, orient=tk.HORIZONTAL, bg=BTN_BG, troughcolor=BG)
        hsb.pack(side=tk.BOTTOM, fill=tk.X)

        text = tk.Text(frame, bg="#1e1e1e", fg=FG, font=("Consolas", self._text_font_size),
                       wrap=tk.NONE, relief=tk.FLAT, padx=12, pady=8,
                       yscrollcommand=vsb.set, xscrollcommand=hsb.set,
                       insertbackground=FG, selectbackground="#264f78")
        text.pack(fill=tk.BOTH, expand=True)
        vsb.config(command=text.yview)
        hsb.config(command=text.xview)

        text.insert(tk.END, content)
        text.config(state=tk.DISABLED)

        # JSON 語法標色（簡易）
        ext = os.path.splitext(filepath)[1].lower()
        if ext == '.json':
            text.config(state=tk.NORMAL)
            self._highlight_json(text)
            text.config(state=tk.DISABLED)

        # 記錄大預覽的 Text widget，並在視窗關閉時移除
        self._open_text_viewers.append(text)
        # 初始化此視窗的字體大小追蹤
        self._viewer_font_sizes[text] = self._text_font_size

        def _on_close_view(_ev=None, tw=text):
            try:
                if tw in self._open_text_viewers:
                    self._open_text_viewers.remove(tw)
                if tw in self._viewer_font_sizes:
                    del self._viewer_font_sizes[tw]
                if self._last_active_text_widget is tw:
                    self._last_active_text_widget = self._text_widget
                # Unbind window-level and widget-level handlers to avoid duplicates
                try:
                    win.unbind("<Control-MouseWheel>")
                except Exception:
                    pass
                try:
                    win.unbind("<space>")
                except Exception:
                    pass
                try:
                    tw.unbind("<Button-1>")
                except Exception:
                    pass
            except Exception:
                pass

        win.bind("<Destroy>", _on_close_view)

        # Ctrl + 滾輪：針對此大預覽視窗調整字體大小（只影響該視窗）
        def _ctrl_wheel(ev, tw=text):
            if ev.delta > 0:
                self._adjust_font_for_widget(tw, 1)
            else:
                self._adjust_font_for_widget(tw, -1)
            return "break"

        # 綁定到大視窗，以便在該視窗聚焦時使用
        win.bind("<Control-MouseWheel>", _ctrl_wheel)
        # 點擊視窗中的文字時，把該視窗設為最後活動目標
        text.bind("<Button-1>", lambda e, tw=text: self._set_active_text_widget(tw))

        # 搜尋功能
        def do_search(*_):
            # remove previous highlights
            text.tag_remove("search_hl", "1.0", tk.END)
            text.tag_remove("search_current", "1.0", tk.END)
            kw = search_var.get()
            if not kw:
                match_var.set("")
                # reset search index state
                try:
                    if text in self._viewer_search_index:
                        del self._viewer_search_index[text]
                except Exception:
                    pass
                return
            count = 0
            start = "1.0"
            text.config(state=tk.NORMAL)
            while True:
                pos = text.search(kw, start, stopindex=tk.END, nocase=True)
                if not pos:
                    break
                end = f"{pos}+{len(kw)}c"
                text.tag_add("search_hl", pos, end)
                start = end
                count += 1
            text.tag_config("search_hl", background="#ff8c00", foreground="#000000")
            # current-match tag (stronger highlight)
            text.tag_config("search_current", background="#ffd700", foreground="#000000")
            text.config(state=tk.DISABLED)
            match_var.set(f"{count} 個符合" if count else "無符合")
            if count:
                # initialize per-viewer index and show first match
                try:
                    self._viewer_search_index[text] = -1
                except Exception:
                    pass
                self._goto_next_search(text)

        search_entry.bind("<Return>", do_search)
        tk.Button(toolbar, text="搜尋", command=do_search,
                  bg=ACCENT, fg="white", relief=tk.FLAT, font=FONT,
                  padx=8, cursor="hand2").pack(side=tk.LEFT)

        # Space => jump to next search result (when this big viewer window is focused)
        def _space_next(ev, tw=text):
            # ignore if no matches
            self._goto_next_search(tw)
            return "break"

        win.bind("<space>", _space_next)

    def _adjust_font_for_widget(self, tw, delta):
        """Adjust font size for a specific Text widget (used by big preview controls)."""
        cur = self._viewer_font_sizes.get(tw, self._text_font_size)
        new = max(7, min(28, cur + delta))
        if new == cur:
            return
        self._viewer_font_sizes[tw] = new
        try:
            tw.config(font=("Consolas", new))
        except Exception:
            pass

    def _restore_ui_control_fonts(self):
        """Force label-type UI controls to keep the app's UI font.
        Deliberately does NOT modify Button widgets — changing button font
        causes Tk to recalculate their pixel size, which triggers pack reflow
        and makes the A+/A− buttons appear to jump or disappear."""
        try:
            frames = []
            for name in ("media_bar", "_pdf_bar", "_text_bar", "_sidebar_frame"):
                f = getattr(self, name, None)
                if isinstance(f, tk.Frame):
                    frames.append(f)

            for fr in frames:
                for w in fr.winfo_children():
                    try:
                        if isinstance(w, (tk.Label, tk.Checkbutton, tk.Radiobutton)):
                            w.config(font=FONT)
                    except Exception:
                        pass
        except Exception:
            pass

    @staticmethod
    def _highlight_json(text_widget):
        import re
        rules = [
            (r'"(\\.|[^"\\])*"\s*:', "#9cdcfe"),   # key
            (r':\s*"(\\.|[^"\\])*"',  "#ce9178"),   # string value
            (r'\b(true|false|null)\b', "#569cd6"),   # keyword
            (r'\b-?\d+(\.\d+)?\b',    "#b5cea8"),   # number
        ]
        for pattern, color in rules:
            tag = f"json_{color}"
            text_widget.tag_config(tag, foreground=color)
            start = "1.0"
            while True:
                match = text_widget.search(pattern, start, stopindex=tk.END, regexp=True)
                if not match:
                    break
                _, _ = match.split('.')
                raw = text_widget.get(match, f"{match} lineend")
                m = re.match(pattern, raw)
                if not m:
                    start = f"{match}+1c"
                    continue
                end = f"{match}+{len(m.group(0))}c"
                text_widget.tag_add(tag, match, end)
                start = end

    def _goto_next_search(self, text_widget):
        """Move to the next search highlight in the given Text widget (cycles)."""
        try:
            ranges = text_widget.tag_ranges("search_hl")
            if not ranges:
                return
            # get list of start indices
            starts = [ranges[i] for i in range(0, len(ranges), 2)]
            cur_idx = self._viewer_search_index.get(text_widget, -1)
            cur_idx = (cur_idx + 1) % len(starts)
            # update stored index
            self._viewer_search_index[text_widget] = cur_idx
            start = starts[cur_idx]
            # corresponding end is the paired range entry
            end = ranges[cur_idx * 2 + 1]
        except Exception:
            # fallback: derive end from first matching tag range
            try:
                ranges = text_widget.tag_ranges("search_hl")
                if not ranges:
                    return
                start = ranges[0]
                end = ranges[1]
            except Exception:
                return

        try:
            # remove previous current tag, add to current match
            text_widget.tag_remove("search_current", "1.0", tk.END)
            text_widget.tag_add("search_current", start, end)
            # ensure visible and set insertion to start
            text_widget.see(start)
            # briefly flash selection by setting insert mark
            try:
                text_widget.mark_set("insert", start)
            except Exception:
                pass
        except Exception:
            pass

    # ── 媒體播放控制 ──────────────────────────────────────────────────────

    def _play(self):
        if not self.files:
            return
        filepath = self.files[self.current_index]
        ext      = os.path.splitext(filepath)[1].lower()

        if self._paused:
            # 從暫停恢復
            self._paused = False
            if ext in VIDEO_EXTS and VLC_AVAILABLE and self._vlc_player:
                self._vlc_player.set_pause(0)
            elif ext in AUDIO_EXTS and PYGAME_AVAILABLE:
                pygame.mixer.music.unpause()
            # cv2 影片：worker thread 偵測 _paused flag 自動繼續
            self.media_status_var.set("▶ 播放中")
            self.play_btn.config(state=tk.DISABLED)
            self.pause_btn.config(state=tk.NORMAL)
            return

        self._stop()

        if ext in AUDIO_EXTS:
            self._play_audio(filepath)
        elif ext in VIDEO_EXTS:
            self._play_video(filepath)

    def _play_audio(self, filepath):
        if not PYGAME_AVAILABLE:
            messagebox.showinfo("提示", "需安裝 pygame:\npip install pygame")
            return
        try:
            pygame.mixer.music.load(filepath)
            pygame.mixer.music.play()
            self._playing = True
            self._paused  = False
            self.media_status_var.set("▶ 播放中")
            self.play_btn.config(state=tk.DISABLED)
            self.pause_btn.config(state=tk.NORMAL)
            self.stop_btn.config(state=tk.NORMAL)
            self._poll_audio_end()
        except Exception as e:
            messagebox.showerror("錯誤", f"無法播放音訊:\n{e}")

    def _poll_audio_end(self):
        if not self._playing:
            return
        if PYGAME_AVAILABLE and not pygame.mixer.music.get_busy() and not self._paused:
            self._stop()
            self.media_status_var.set("⏹ 播放完畢")
            return
        self._after_id = self.root.after(300, self._poll_audio_end)

    def _play_video(self, filepath):
        if VLC_AVAILABLE:
            self._play_video_vlc(filepath)
        elif CV2_AVAILABLE and PIL_AVAILABLE:
            self._play_video_cv2(filepath)
        else:
            messagebox.showinfo("提示",
                "需安裝 python-vlc（含聲音）:\n  pip install python-vlc\n\n"
                "或安裝 opencv+Pillow（靜音）:\n  pip install opencv-python Pillow")

    def _play_video_vlc(self, filepath):
        # 延遲初始化 VLC instance
        if self._vlc_instance is None:
            # Disable hardware-accelerated decoding to avoid driver/D3D11 issues
            # (some systems produce errors like "SetThumbNailClip failed" when using hwdec)
            try:
                self._vlc_instance = vlc.Instance(['--no-video-title-show', '--avcodec-hw=none'])
            except Exception:
                # fallback to default instance if passing options fails
                self._vlc_instance = vlc.Instance()
            self._vlc_player   = self._vlc_instance.media_player_new()

        media = self._vlc_instance.media_new(filepath)
        self._vlc_player.set_media(media)

        # 切到 VLC 專用 Frame，嘗試綁定 HWND（若失敗則回退到縮圖顯示）
        self._switch_to_vlc_view()
        self.root.update_idletasks()
        try:
            self._vlc_player.set_hwnd(self._vlc_frame.winfo_id())
        except Exception:
            # 如果無法綁定 HWND（Direct3D/驅動問題），回退到 canvas 顯示縮圖並通知使用者
            try:
                self._switch_to_canvas_view()
                self._show_video_thumb(filepath)
            except Exception:
                self._set_canvas_text("影片無法直接播放（VLC HWND 綁定失敗）")
            messagebox.showwarning("提示", "VLC 無法綁定視窗，已使用靜態縮圖代替播放（硬體加速可能導致錯誤）。")
            return

        self._vlc_playing_path = filepath   # 記錄實際播放的路徑，停止時用
        # reset fallback flag (to avoid infinite loops)
        self._vlc_fallback_tried = False
        self._vlc_player.play()
        self._playing = True
        self._paused  = False
        self.media_status_var.set("▶ 播放中（含聲音）")
        self.play_btn.config(state=tk.DISABLED)
        self.pause_btn.config(state=tk.NORMAL)
        self.stop_btn.config(state=tk.NORMAL)
        self._poll_vlc_end()

    def _poll_vlc_end(self):
        if not self._playing:
            return
        state = self._vlc_player.get_state()
        if state in (vlc.State.Ended, vlc.State.Stopped):
            self._stop()
            self.media_status_var.set("⏹ 播放完畢")
            return
        if state == vlc.State.Error:
            # VLC encountered an error (often hwdec / D3D11 issues). Try fallback to cv2 playback if available.
            if not getattr(self, '_vlc_fallback_tried', False) and CV2_AVAILABLE and PIL_AVAILABLE:
                self._vlc_fallback_tried = True
                try:
                    self._stop()
                except Exception:
                    pass
                messagebox.showwarning("播放錯誤", "VLC 播放發生錯誤，已改用靜音模式播放（OpenCV）。")
                try:
                    self._play_video_cv2(self._vlc_playing_path)
                except Exception:
                    self._set_canvas_text("無法播放影片（VLC 與 OpenCV 都失敗）")
                return
            else:
                self._stop()
                self.media_status_var.set("⏹ 播放失敗")
                return
        self._after_id = self.root.after(500, self._poll_vlc_end)

    def _play_video_cv2(self, filepath):
        self._stop_event.clear()
        self._playing = True
        self._paused  = False
        self.media_status_var.set("▶ 播放中（靜音）")
        self.play_btn.config(state=tk.DISABLED)
        self.pause_btn.config(state=tk.NORMAL)
        self.stop_btn.config(state=tk.NORMAL)

        self._video_thread = threading.Thread(
            target=self._video_worker, args=(filepath,), daemon=True)
        self._video_thread.start()
        self._video_frame_loop()

    def _video_worker(self, filepath):
        cap = cv2.VideoCapture(filepath)
        fps = cap.get(cv2.CAP_PROP_FPS) or 30
        delay = 1.0 / fps

        while not self._stop_event.is_set():
            if self._paused:
                time.sleep(0.05)
                continue
            ret, frame = cap.read()
            if not ret:
                break
            try:
                self._frame_q.put(frame, timeout=0.1)
            except queue.Full:
                pass
            time.sleep(delay)

        cap.release()
        self._frame_q.put(None)  # 結束信號

    def _video_frame_loop(self):
        if not self._playing:
            return
        try:
            frame = self._frame_q.get_nowait()
            if frame is None:
                # 播放結束
                self._stop()
                self.media_status_var.set("⏹ 播放完畢")
                return
            self._render_cv2_frame(frame)
        except queue.Empty:
            pass
        self._after_id = self.root.after(16, self._video_frame_loop)

    def _pause(self):
        if not self._playing:
            return
        filepath = self.files[self.current_index]
        ext      = os.path.splitext(filepath)[1].lower()

        self._paused = True
        if ext in VIDEO_EXTS and VLC_AVAILABLE and self._vlc_player:
            self._vlc_player.set_pause(1)
        elif ext in AUDIO_EXTS and PYGAME_AVAILABLE:
            pygame.mixer.music.pause()
        self.media_status_var.set("⏸ 已暫停")
        self.play_btn.config(state=tk.NORMAL, text="▶  繼續")
        self.pause_btn.config(state=tk.DISABLED)

    def _toggle_pause(self):
        if self._paused:
            self._play()
        elif self._playing:
            self._pause()

    def _stop(self):
        if self._after_id:
            self.root.after_cancel(self._after_id)
            self._after_id = None

        self._stop_event.set()
        self._playing = False
        self._paused  = False

        # 等 cv2 執行緒結束（最多 0.5 秒，避免殘幀問題）
        if self._video_thread and self._video_thread.is_alive():
            self._video_thread.join(timeout=0.5)
        self._video_thread = None

        # 清空 frame queue
        while not self._frame_q.empty():
            try:
                self._frame_q.get_nowait()
            except queue.Empty:
                break

        if VLC_AVAILABLE and self._vlc_player:
            try:
                # 關鍵順序：set_hwnd(0) 必須在 stop() 之前
                # 這樣 VLC 渲染執行緒立即失去目標，不再碰 widget 的 HWND
                self._vlc_player.set_hwnd(0)
                self._vlc_player.stop()
            except Exception:
                pass
            # HWND 已脫鉤，現在才安全地操作 widget 佈局
            if self._vlc_frame.winfo_manager():
                self._switch_to_canvas_view()
                # 恢復靜態縮圖：使用 VLC 播放時記錄的路徑，
                # 而非 current_index（用戶可能已導航到其他檔案）
                fp = getattr(self, '_vlc_playing_path', None)
                if fp and os.path.splitext(fp)[1].lower() in VIDEO_EXTS:
                    try:
                        self._show_video_thumb(fp)
                    except Exception:
                        pass

        if PYGAME_AVAILABLE:
            try:
                pygame.mixer.music.stop()
            except Exception:
                pass

        self.play_btn.config(state=tk.NORMAL, text="▶  播放")
        self.pause_btn.config(state=tk.DISABLED)
        self.stop_btn.config(state=tk.DISABLED)
        self.media_status_var.set("⏹ 已停止")

    # ── 導航 ──────────────────────────────────────────────────────────────

    def _prev(self):
        if self.files and self.current_index > 0:
            self.current_index -= 1
            self._thumb_update_highlight()
            self._show_current()

    def _next(self):
        if self.files and self.current_index < len(self.files) - 1:
            self.current_index += 1
            self._thumb_update_highlight()
            self._show_current()

    # ── 重新命名 ──────────────────────────────────────────────────────────

    def _rename(self):
        if not self.files:
            return
        new_stem = self.rename_var.get().strip()
        if not new_stem:
            messagebox.showwarning("警告", "請輸入新的檔名（不含副檔名）")
            return

        # 安全性：移除路徑分隔符，防止路徑遍歷攻擊
        new_stem = os.path.basename(new_stem)
        # 過濾 Windows 不允許的檔名字元
        invalid_chars = r'\/:*?"<>|'
        if any(c in new_stem for c in invalid_chars):
            messagebox.showwarning("警告", f"檔名不可包含以下字元：{invalid_chars}")
            return
        if not new_stem:
            messagebox.showwarning("警告", "請輸入有效的檔名")
            return

        old_path = self.files[self.current_index]
        ext      = os.path.splitext(old_path)[1]
        # 保留原始目錄（支援遞迴掃描的子目錄檔案）
        new_path = os.path.join(os.path.dirname(old_path), new_stem + ext)

        if os.path.normcase(old_path) == os.path.normcase(new_path):
            messagebox.showinfo("提示", "檔名未變更")
            return

        if os.path.exists(new_path):
            if not messagebox.askyesno("確認覆蓋", f"「{new_stem + ext}」已存在，要覆蓋嗎?"):
                return

        # 重命名前先停止播放（避免檔案鎖定）
        self._stop()

        try:
            # Use os.replace to ensure overwrite when user confirmed; on Windows os.rename
            # may fail when target exists. os.replace will atomically replace the target.
            os.replace(old_path, new_path)
            self.files[self.current_index] = new_path
            self.info_var.set(f"已重新命名 → {new_stem + ext}")
        except Exception as e:
            messagebox.showerror("錯誤", f"重新命名失敗:\n{e}")

    # ── 縮圖側欄 ──────────────────────────────────────────────────────────

    def _sidebar_rebuild(self):
        """重建縮圖側欄：繪製佔位格子，啟動背景載入執行緒。"""
        self._thumb_gen += 1
        self._thumb_stop.set()
        if self._thumb_thread and self._thumb_thread.is_alive():
            self._thumb_thread.join(timeout=0.3)
        self._thumb_stop.clear()

        self._thumb_canvas.delete("all")
        self._thumb_slots = []

        y = 4
        for i, fp in enumerate(self.files):
            fname = os.path.basename(fp)
            if len(fname) > 15:
                fname = fname[:13] + "…"
            bg_r = self._thumb_canvas.create_rectangle(
                2, y, 116, y + THUMB_ITEM_H - 2,
                fill="#252526", outline="#3c3c3c", width=1, tags=(f"slot_{i}",))
            ph_r = self._thumb_canvas.create_rectangle(
                4, y + 2, 114, y + THUMB_H + 2,
                fill="#3c3c3c", outline="", tags=(f"ph_{i}",))
            lbl  = self._thumb_canvas.create_text(
                59, y + THUMB_H + 11, text=fname,
                fill="#9cdcfe", font=("Segoe UI", 7),
                width=108, tags=(f"lbl_{i}",))
            for item in (bg_r, ph_r, lbl):
                self._thumb_canvas.tag_bind(
                    item, "<Button-1>",
                    lambda e, idx=i: self._thumb_click(idx))
                self._thumb_canvas.tag_bind(
                    item, "<Double-1>",
                    lambda e, idx=i: self._thumb_dblclick(idx))
            self._thumb_slots.append({"bg": bg_r, "ph": ph_r, "lbl": lbl, "img": None})
            y += THUMB_ITEM_H

        self._thumb_canvas.configure(scrollregion=(0, 0, 118, y + 4))
        self._thumb_update_highlight()

        gen = self._thumb_gen
        self._thumb_thread = threading.Thread(
            target=self._thumb_worker, args=(list(self.files), gen), daemon=True)
        self._thumb_thread.start()

    def _thumb_worker(self, files, gen):
        for i, fp in enumerate(files):
            if self._thumb_stop.is_set() or gen != self._thumb_gen:
                break
            if fp in self._thumb_photos:
                p = self._thumb_photos[fp]
                self.root.after(0, lambda idx=i, ph=p, g=gen: self._thumb_apply(idx, ph, g))
                time.sleep(0.005)
                continue
            img = self._thumb_make(fp)
            if img is not None:
                self.root.after(0, lambda idx=i, im=img, path=fp, g=gen:
                                self._thumb_to_photo(idx, im, path, g))
            time.sleep(0.01)

    def _thumb_make(self, fp):
        ext = os.path.splitext(fp)[1].lower()
        try:
            if ext in IMAGE_EXTS and PIL_AVAILABLE:
                img = Image.open(fp)
                img.thumbnail((THUMB_W, THUMB_H), Image.LANCZOS)
                return img
            if ext in VIDEO_EXTS and CV2_AVAILABLE and PIL_AVAILABLE:
                cap = cv2.VideoCapture(fp)
                ret, frame = cap.read()
                cap.release()
                if ret:
                    img = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
                    img.thumbnail((THUMB_W, THUMB_H), Image.LANCZOS)
                    return img
            if ext in PDF_EXTS and FITZ_AVAILABLE and PIL_AVAILABLE:
                doc = _fitz.open(fp)
                # 若需密碼，先驗證
                if doc.needs_pass:
                    if fp in self._pdf_passwords:
                        pwd = self._pdf_passwords[fp]
                        auth_code = doc.authenticate(pwd)
                        if auth_code == 0:
                            doc.close()
                            return None  # 密碼無效
                    else:
                        doc.close()
                        return None  # 需要密碼但未提供
                pix = doc[0].get_pixmap(matrix=_fitz.Matrix(0.22, 0.22))
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                doc.close()
                img.thumbnail((THUMB_W, THUMB_H), Image.LANCZOS)
                return img
            if ext in EXCEL_EXTS and PIL_AVAILABLE:
                # 綠底 mini 表格縮圖
                img  = Image.new('RGB', (THUMB_W, THUMB_H), '#1a2a1a')
                draw = ImageDraw.Draw(img)
                # Excel 綠色標題列
                draw.rectangle([0, 0, THUMB_W - 1, 13], fill='#217346')
                draw.text((3, 2), 'XLSX', fill='#ffffff')
                # 讀取前 5 列 × 4 欄顯示成格線
                if OPENPYXL_AVAILABLE:
                    try:
                        wb = _openpyxl.load_workbook(fp, read_only=True, data_only=True)
                        ws = wb.active
                        cell_rows = []
                        for row in ws.iter_rows(min_row=1, max_row=5,
                                                max_col=4, values_only=True):
                            cell_rows.append(row)
                        wb.close()
                        col_w = max(1, THUMB_W // 4)
                        row_h = max(1, (THUMB_H - 15) // max(len(cell_rows), 1))
                        for ri, row in enumerate(cell_rows):
                            for ci, val in enumerate(row):
                                x1 = ci * col_w
                                y1 = 15 + ri * row_h
                                draw.rectangle(
                                    [x1, y1, x1 + col_w - 1, y1 + row_h - 1],
                                    outline='#2d6a3a', fill='#1e3a22')
                                if val is not None:
                                    draw.text((x1 + 2, y1 + 1),
                                              str(val)[:6], fill='#9cdcfe')
                    except Exception:
                        pass
                return img
        except Exception:
            pass
        return None

    def _thumb_to_photo(self, idx, img, fp, gen):
        if gen != self._thumb_gen:
            return
        photo = ImageTk.PhotoImage(img)
        # 縮圖快取上限 200 張，避免記憶體無限成長
        if len(self._thumb_photos) >= 200:
            try:
                oldest_key = next(iter(self._thumb_photos))
                del self._thumb_photos[oldest_key]
            except StopIteration:
                pass
        self._thumb_photos[fp] = photo
        self._thumb_apply(idx, photo, gen)

    def _thumb_apply(self, idx, photo, gen):
        if gen != self._thumb_gen or idx >= len(self._thumb_slots):
            return
        slot = self._thumb_slots[idx]
        y = 4 + idx * THUMB_ITEM_H
        self._thumb_canvas.delete(slot["ph"])
        img_item = self._thumb_canvas.create_image(
            59, y + 2 + THUMB_H // 2, image=photo, tags=(f"img_{idx}",))
        self._thumb_canvas.tag_bind(
            img_item, "<Button-1>",
            lambda e, i=idx: self._thumb_click(i))
        self._thumb_canvas.tag_bind(
            img_item, "<Double-1>",
            lambda e, i=idx: self._thumb_dblclick(i))
        self._thumb_canvas.tag_raise(f"lbl_{idx}")
        slot["img"] = img_item

    def _thumb_click(self, idx):
        self.current_index = idx
        self._thumb_update_highlight()
        self._show_current()

    def _thumb_dblclick(self, idx):
        """雙擊縮圖：用系統預設程式直接開啟該檔案。"""
        if 0 <= idx < len(self.files):
            self._open_with_system(self.files[idx])

    def _thumb_update_highlight(self):
        for i, slot in enumerate(self._thumb_slots):
            selected = (i == self.current_index)
            # 選中：亮藍填色 + 明顯邊框；未選中：暗底無邊框
            self._thumb_canvas.itemconfig(
                slot["bg"],
                fill    = "#1a3a5c" if selected else "#252526",
                outline = "#3a9ddc" if selected else "#3c3c3c",
                width   = 2        if selected else 1)
            # 標籤文字顏色跟著變
            self._thumb_canvas.itemconfig(
                slot["lbl"],
                fill = "#ffffff" if selected else "#9cdcfe")

        # 確保選中項可見（捲動到對應位置）
        if self._thumb_slots and 0 <= self.current_index < len(self._thumb_slots):
            n     = len(self._thumb_slots)
            total = n * THUMB_ITEM_H + 8
            # 讓選中項的頂端出現在視窗上方 1/4 處（向上偏移讓它置中可見）
            item_top = 4 + self.current_index * THUMB_ITEM_H
            frac = max(0.0, (item_top - 30) / max(total, 1))
            self._thumb_canvas.yview_moveto(frac)

    def _thumb_on_wheel(self, event):
        self._thumb_canvas.yview_scroll(-1 if event.delta > 0 else 1, "units")

    # ── 拖曳載入（Drag & Drop）──────────────────────────────────────────────

    def _setup_dnd(self):
        # 不管是否有 DnD，都顯示初始提示
        self._set_canvas_text(
            "尚未載入檔案\n\n"
            "① 上方輸入路徑 → 載入\n"
            "② 點「資料夾」或「單一檔案」瀏覽\n"
            + ("③ 直接拖曳檔案 / 資料夾到此處" if DND_AVAILABLE
               else "③ 安裝 tkinterdnd2 可啟用拖曳功能"))

        if not DND_AVAILABLE:
            return
        try:
            # 根視窗與 canvas 都登記為拖放目標
            self.root.drop_target_register(DND_FILES)
            self.root.dnd_bind("<<Drop>>",      self._on_drop)
            self.root.dnd_bind("<<DragEnter>>",  self._on_drag_enter)
            self.root.dnd_bind("<<DragLeave>>",  self._on_drag_leave)
        except Exception:
            pass

    def _on_drag_enter(self, _event=None):
        """拖曳進入時：僅加亮邊框；有檔案時保留現有預覽不覆蓋。"""
        self.canvas.config(highlightthickness=3, highlightbackground="#3a9ddc")
        if not self.files:
            self._set_canvas_text("放開滑鼠以載入檔案 ↓")

    def _on_drag_leave(self, _event=None):
        """拖曳離開時：還原 canvas 外觀。"""
        self.canvas.config(highlightthickness=0)
        if not self.files:
            self._set_canvas_text(
                "尚未載入檔案\n\n"
                "① 上方輸入路徑 → 載入\n"
                "② 點「資料夾」或「單一檔案」瀏覽\n"
                "③ 直接拖曳檔案 / 資料夾到此處")

    def _on_drop(self, event):
        if not self._ensure_stopped_for_action():
            return
        try:
            raw = self.root.tk.splitlist(event.data)
        except Exception:
            raw = event.data.split()
        all_files = []
        primary   = None
        for path in raw:
            path = path.strip().strip('"')   # 清理引號與空白（Windows DnD 常帶引號）
            if not path:
                continue
            # 正規化路徑，防止路徑遍歷
            path = os.path.normpath(path)
            if os.path.isdir(path):
                if primary is None:
                    primary = path
                for dp, _, fnames in os.walk(path):
                    for fn in fnames:
                        if os.path.splitext(fn)[1].lower() in SUPPORTED_EXTS:
                            all_files.append(os.path.join(dp, fn))
            elif os.path.isfile(path):
                if os.path.splitext(path)[1].lower() in SUPPORTED_EXTS:
                    all_files.append(path)
        self.canvas.config(highlightthickness=0)   # 移除拖曳高亮
        if not all_files:
            self._on_drag_leave()   # 還原提示
            return
        self._stop()
        self._pdf_cancel_all()
        self.folder_path   = primary or os.path.dirname(all_files[0])
        self.files         = sorted(set(all_files))
        self.current_index = 0
        self.path_var.set(self.folder_path if primary else all_files[0])
        self._sidebar_rebuild()
        self._show_current()

    # ── 全域搜尋 ──────────────────────────────────────────────────────────

    def _open_search(self):
        if not self._ensure_stopped_for_action():
            return
        SearchWindow(self.root, self.folder_path, self.files, self._navigate_to_file)

    def _navigate_to_file(self, filepath):
        """從搜尋結果導航到指定檔案。"""
        if filepath in self.files:
            self.current_index = self.files.index(filepath)
        else:
            self._stop()
            self.files.append(filepath)
            self.files.sort()
            self.current_index = self.files.index(filepath)
            self.folder_path = os.path.dirname(filepath)
            self._sidebar_rebuild()
        self._thumb_update_highlight()
        self._show_current()

    # ── 批次重命名 ────────────────────────────────────────────────────────

    def _open_batch_rename(self):
        if not self.files:
            messagebox.showinfo("提示", "請先載入資料夾")
            return
        if not self._ensure_stopped_for_action():
            return
        def on_done():
            self._load_files()
        BatchRenameDialog(self.root, list(self.files), self.folder_path, on_done)

    # ── 關閉視窗 ──────────────────────────────────────────────────────────

    def _on_close(self):
        self._stop()
        self._thumb_stop.set()
        self._pdf_cancel_all()
        if self._pdf_doc:
            try:
                self._pdf_doc.close()
            except Exception:
                pass
        if VLC_AVAILABLE and self._vlc_player:
            try:
                self._vlc_player.release()
            except Exception:
                pass
        if VLC_AVAILABLE and self._vlc_instance:
            try:
                self._vlc_instance.release()
            except Exception:
                pass
        if PYGAME_AVAILABLE:
            pygame.mixer.quit()
        self.root.destroy()


# ── 全域搜尋視窗 ──────────────────────────────────────────────────────────

class SearchWindow:
    def __init__(self, parent, root_path, files, on_navigate):
        self.win = tk.Toplevel(parent)
        self.win.title("全域搜尋")
        self.win.geometry("940x620")
        self.win.configure(bg=BG)
        self.win.resizable(True, True)

        self.root_path   = root_path
        self.files       = list(files)
        self.on_navigate = on_navigate
        self._stop_flag  = threading.Event()
        self._results    = []
        self._build()

    def _build(self):
        # ── 控制列 ──
        ctrl = tk.Frame(self.win, bg=BG, pady=8)
        ctrl.pack(fill=tk.X, padx=12)

        tk.Label(ctrl, text="關鍵字:", bg=BG, fg=FG, font=FONT).pack(side=tk.LEFT)
        self._kw_var = tk.StringVar()
        kw_e = tk.Entry(ctrl, textvariable=self._kw_var, bg=ENTRY_BG, fg=FG,
                        insertbackground=FG, relief=tk.FLAT, font=FONT, width=26)
        kw_e.pack(side=tk.LEFT, padx=6, ipady=4)
        kw_e.bind("<Return>", lambda _: self._start())
        kw_e.focus_set()

        self._case_var  = tk.BooleanVar(value=False)
        self._regex_var = tk.BooleanVar(value=False)
        self._scope_var = tk.StringVar(value="folder")

        tk.Checkbutton(ctrl, text="大小寫", variable=self._case_var,
                       bg=BG, fg=FG, selectcolor=BTN_BG, activebackground=BG,
                       font=FONT).pack(side=tk.LEFT, padx=3)
        tk.Checkbutton(ctrl, text="Regex", variable=self._regex_var,
                       bg=BG, fg=FG, selectcolor=BTN_BG, activebackground=BG,
                       font=FONT).pack(side=tk.LEFT)
        tk.Frame(ctrl, bg=BTN_BG, width=1).pack(side=tk.LEFT, fill=tk.Y, padx=8)
        tk.Radiobutton(ctrl, text="資料夾", variable=self._scope_var, value="folder",
                       bg=BG, fg=FG, selectcolor=BTN_BG, activebackground=BG,
                       font=FONT).pack(side=tk.LEFT)
        tk.Radiobutton(ctrl, text="目前清單", variable=self._scope_var, value="list",
                       bg=BG, fg=FG, selectcolor=BTN_BG, activebackground=BG,
                       font=FONT).pack(side=tk.LEFT, padx=4)

        tk.Button(ctrl, text="搜尋", command=self._start,
                  bg=ACCENT, fg="white", relief=tk.FLAT, font=FONT,
                  padx=10, cursor="hand2").pack(side=tk.LEFT, padx=6)
        tk.Button(ctrl, text="停止", command=self._stop_search,
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT,
                  padx=8, cursor="hand2").pack(side=tk.LEFT)

        tk.Button(ctrl, text="匯出 JSON", command=lambda: self._export("json"),
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT,
                  padx=6, cursor="hand2").pack(side=tk.RIGHT)
        tk.Button(ctrl, text="匯出 TXT",  command=lambda: self._export("txt"),
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT,
                  padx=6, cursor="hand2").pack(side=tk.RIGHT, padx=4)

        self._status_var = tk.StringVar(value="請輸入關鍵字並按搜尋")
        tk.Label(ctrl, textvariable=self._status_var, bg=BG, fg="#9cdcfe",
                 font=("Segoe UI", 9)).pack(side=tk.LEFT, padx=10)

        # ── 進度條 ──
        self._progress = ttk.Progressbar(self.win, mode="indeterminate", length=100)
        self._progress.pack(fill=tk.X, padx=12, pady=(0, 4))

        # ── 結果 Treeview ──
        tf = tk.Frame(self.win, bg="#1a1a1a")
        tf.pack(fill=tk.BOTH, expand=True, padx=12, pady=(0, 8))

        vsb = tk.Scrollbar(tf, bg=BTN_BG, troughcolor=BG)
        vsb.pack(side=tk.RIGHT, fill=tk.Y)
        hsb = tk.Scrollbar(tf, orient=tk.HORIZONTAL, bg=BTN_BG, troughcolor=BG)
        hsb.pack(side=tk.BOTTOM, fill=tk.X)

        style = ttk.Style()
        style.theme_use("default")
        style.configure("S.Treeview", background="#1e1e1e", foreground=FG,
                        fieldbackground="#1e1e1e", rowheight=22)
        style.configure("S.Treeview.Heading", background="#252526",
                        foreground=FG, relief="flat")
        style.map("S.Treeview", background=[("selected", "#264f78")])

        self._tree = ttk.Treeview(tf, columns=("file", "line", "match", "folder"),
                                   show="headings", style="S.Treeview",
                                   yscrollcommand=vsb.set, xscrollcommand=hsb.set)
        self._tree.heading("file",   text="檔案名稱")
        self._tree.heading("line",   text="行號")
        self._tree.heading("match",  text="匹配內容")
        self._tree.heading("folder", text="所在資料夾")
        self._tree.column("file",   width=160, minwidth=60)
        self._tree.column("line",   width=55,  minwidth=40, anchor="center")
        self._tree.column("match",  width=450, minwidth=100)
        self._tree.column("folder", width=220, minwidth=60)
        self._tree.pack(fill=tk.BOTH, expand=True)
        vsb.config(command=self._tree.yview)
        hsb.config(command=self._tree.xview)
        self._tree.bind("<Double-1>", self._on_dclick)
        self._tree.bind("<Return>",   self._on_dclick)

    # ── 搜尋邏輯 ──────────────────────────────────────────────────────────

    def _target_files(self):
        if self._scope_var.get() == "list":
            return [f for f in self.files
                    if os.path.splitext(f)[1].lower() in SEARCH_EXTS]
        found = []
        base = self.root_path or ""
        if os.path.isdir(base):
            for dp, _, fnames in os.walk(base):
                for fn in fnames:
                    if os.path.splitext(fn)[1].lower() in SEARCH_EXTS:
                        found.append(os.path.join(dp, fn))
        return found

    def _start(self):
        kw = self._kw_var.get().strip()
        if not kw:
            return
        self._stop_flag.clear()
        self._tree.delete(*self._tree.get_children())
        self._results.clear()
        self._progress.start(12)
        self._status_var.set("搜尋中…")
        threading.Thread(
            target=self._worker,
            args=(kw, self._target_files(),
                  self._regex_var.get(), self._case_var.get()),
            daemon=True).start()

    @staticmethod
    def _extract_lines(fp):
        """將任意支援格式轉成 [(lineno, text), ...] 供搜尋用。"""
        ext = os.path.splitext(fp)[1].lower()
        try:
            if ext == '.docx':
                if not DOCX_AVAILABLE:
                    return []
                doc = _docx.Document(fp)
                return [(i + 1, p.text) for i, p in enumerate(doc.paragraphs) if p.text]
            if ext == '.pdf':
                if not FITZ_AVAILABLE:
                    return []
                doc = _fitz.open(fp)
                lines = []
                global_line = 1
                for page in doc:
                    for raw in page.get_text().splitlines():
                        if raw.strip():
                            lines.append((global_line, raw))
                        global_line += 1
                doc.close()
                return lines
            if ext in EXCEL_EXTS:
                if not OPENPYXL_AVAILABLE:
                    return []
                wb = _openpyxl.load_workbook(fp, read_only=True, data_only=True)
                lines = []
                global_line = 1
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    for row in ws.iter_rows(values_only=True):
                        text = '\t'.join(str(c) for c in row if c is not None)
                        if text.strip():
                            lines.append((global_line, f"[{sheet_name}] {text}"))
                        global_line += 1
                wb.close()
                return lines
        except Exception:
            return []
        # 純文字類：逐行讀取
        try:
            with open(fp, encoding="utf-8", errors="replace") as f:
                return list(enumerate(f, 1))
        except Exception:
            return []

    def _worker(self, kw, files, use_regex, case_sens):
        import re as _re
        pattern = None
        if use_regex:
            if len(kw) > 200:
                self.win.after(0, lambda: self._status_var.set("Regex 過長（限 200 字元）"))
                self.win.after(0, self._progress.stop)
                return
            flags = 0 if case_sens else _re.IGNORECASE
            try:
                pattern = _re.compile(kw, flags)
            except _re.error:
                self.win.after(0, lambda: self._status_var.set("Regex 語法錯誤"))
                self.win.after(0, self._progress.stop)
                return
        needle = kw if case_sens else kw.lower()
        count  = 0
        for fp in files:
            if self._stop_flag.is_set():
                break
            for lineno, line in self._extract_lines(fp):
                if self._stop_flag.is_set():
                    break
                hit = (bool(pattern.search(line)) if pattern
                       else needle in (line if case_sens else line.lower()))
                if hit:
                    count += 1
                    row = (os.path.basename(fp), lineno,
                           line.strip()[:100], os.path.dirname(fp), fp)
                    self._results.append(row)
                    self.win.after(0, lambda r=row: self._add_row(r))
        self.win.after(0, lambda: self._done(count))

    def _add_row(self, row):
        fname, lineno, preview, folder, fp = row
        self._tree.insert("", tk.END, values=(fname, lineno, preview, folder),
                          tags=(fp,))

    def _done(self, count):
        self._progress.stop()
        self._status_var.set(f"完成：{count} 個符合")

    def _stop_search(self):
        self._stop_flag.set()
        self._progress.stop()
        self._status_var.set("已停止")

    def _on_dclick(self, _=None):
        sel = self._tree.selection()
        if sel:
            tags = self._tree.item(sel[0], "tags")
            if tags:
                self.on_navigate(tags[0])

    def _export(self, fmt):
        if not self._results:
            messagebox.showinfo("提示", "沒有結果可匯出", parent=self.win)
            return
        path = filedialog.asksaveasfilename(
            parent=self.win, defaultextension=f".{fmt}",
            filetypes=[(f"{fmt.upper()}", f"*.{fmt}"), ("全部", "*.*")])
        if not path:
            return
        try:
            if fmt == "txt":
                with open(path, "w", encoding="utf-8") as f:
                    for fname, lineno, preview, folder, fp in self._results:
                        f.write(f"[{folder}\\{fname}] L{lineno}: {preview}\n")
            else:
                data = [{"file": fp, "line": lineno, "match": preview, "folder": folder}
                        for fname, lineno, preview, folder, fp in self._results]
                with open(path, "w", encoding="utf-8") as f:
                    _json.dump(data, f, ensure_ascii=False, indent=2)
            messagebox.showinfo("完成", f"已匯出 {len(self._results)} 筆", parent=self.win)
        except Exception as e:
            messagebox.showerror("錯誤", str(e), parent=self.win)


# ── 批次重命名對話框 ──────────────────────────────────────────────────────

class BatchRenameDialog:
    def __init__(self, parent, files, folder_path, on_done):
        self.dlg         = tk.Toplevel(parent)
        self.dlg.title("批次重命名")
        self.dlg.geometry("780x540")
        self.dlg.configure(bg=BG)
        self.dlg.resizable(True, True)
        self.dlg.grab_set()

        self.files       = files
        self.folder_path = folder_path
        self.on_done     = on_done

        self.check_vars  = [tk.BooleanVar(value=True) for _ in files]
        self.prefix_var  = tk.StringVar(value="file")
        self.start_var   = tk.StringVar(value="1")
        self.preview_vars = []

        self._build()
        self._update_preview()

    def _build(self):
        # ── 控制列：前綴、起始數、全選 ──
        ctrl = tk.Frame(self.dlg, bg=BG, pady=8)
        ctrl.pack(fill=tk.X, padx=12)

        tk.Label(ctrl, text="前綴:", bg=BG, fg=FG, font=FONT).pack(side=tk.LEFT)
        pe = tk.Entry(ctrl, textvariable=self.prefix_var, bg=ENTRY_BG, fg=FG,
                      insertbackground=FG, relief=tk.FLAT, font=FONT, width=18)
        pe.pack(side=tk.LEFT, padx=(4, 12), ipady=3)
        pe.bind("<KeyRelease>", lambda _: self._update_preview())

        tk.Label(ctrl, text="起始數字:", bg=BG, fg=FG, font=FONT).pack(side=tk.LEFT)
        se = tk.Entry(ctrl, textvariable=self.start_var, bg=ENTRY_BG, fg=FG,
                      insertbackground=FG, relief=tk.FLAT, font=FONT, width=5)
        se.pack(side=tk.LEFT, padx=(4, 12), ipady=3)
        se.bind("<KeyRelease>", lambda _: self._update_preview())

        tk.Button(ctrl, text="全選", command=self._select_all,
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT,
                  padx=6, cursor="hand2").pack(side=tk.LEFT, padx=2)
        tk.Button(ctrl, text="全不選", command=self._deselect_all,
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT,
                  padx=6, cursor="hand2").pack(side=tk.LEFT, padx=2)

        self.count_var = tk.StringVar(value="")
        tk.Label(ctrl, textvariable=self.count_var, bg=BG, fg="#9cdcfe",
                 font=("Segoe UI", 9)).pack(side=tk.LEFT, padx=10)

        # ── 表頭 ──
        hdr = tk.Frame(self.dlg, bg="#252526")
        hdr.pack(fill=tk.X, padx=12)
        tk.Label(hdr, text=" ", bg="#252526", fg=FG, width=2).pack(side=tk.LEFT)
        tk.Label(hdr, text="原始檔名", bg="#252526", fg=FG,
                 font=("Segoe UI", 9, "bold"), width=36, anchor="w").pack(side=tk.LEFT)
        tk.Label(hdr, text="重命名後", bg="#252526", fg=FG,
                 font=("Segoe UI", 9, "bold"), anchor="w").pack(side=tk.LEFT, padx=4)

        # ── 可捲動清單 ──
        list_outer = tk.Frame(self.dlg, bg=BG)
        list_outer.pack(fill=tk.BOTH, expand=True, padx=12, pady=2)

        sb = tk.Scrollbar(list_outer, orient=tk.VERTICAL, bg=BTN_BG)
        sb.pack(side=tk.RIGHT, fill=tk.Y)

        self.lc = tk.Canvas(list_outer, bg="#1a1a1a", highlightthickness=0,
                             yscrollcommand=sb.set)
        self.lc.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        sb.config(command=self.lc.yview)

        self.inner = tk.Frame(self.lc, bg="#1a1a1a")
        self._win_id = self.lc.create_window((0, 0), window=self.inner, anchor="nw")

        self.inner.bind("<Configure>", lambda _: self.lc.configure(
            scrollregion=self.lc.bbox("all")))
        self.lc.bind("<Configure>", lambda e: self.lc.itemconfig(
            self._win_id, width=e.width))

        # 滑鼠滾輪僅在 canvas 上作用
        self.lc.bind("<Enter>",  lambda _: self.lc.bind_all("<MouseWheel>", self._on_wheel))
        self.lc.bind("<Leave>",  lambda _: self.lc.unbind_all("<MouseWheel>"))

        # 建立每一列
        for i, fpath in enumerate(self.files):
            fname   = os.path.basename(fpath)
            row_bg  = "#1e1e1e" if i % 2 == 0 else "#252526"
            row     = tk.Frame(self.inner, bg=row_bg)
            row.pack(fill=tk.X)

            tk.Checkbutton(row, variable=self.check_vars[i],
                           command=self._update_preview,
                           bg=row_bg, fg=FG, selectcolor="#3c3c3c",
                           activebackground=row_bg, relief=tk.FLAT,
                           cursor="hand2").pack(side=tk.LEFT, padx=4)

            tk.Label(row, text=fname, bg=row_bg, fg=FG, font=("Segoe UI", 9),
                     width=36, anchor="w").pack(side=tk.LEFT)

            pv = tk.StringVar(value="")
            self.preview_vars.append(pv)
            tk.Label(row, textvariable=pv, bg=row_bg, fg="#4ec9b0",
                     font=("Segoe UI", 9), anchor="w", width=30).pack(side=tk.LEFT, padx=6)

        # ── 底部：執行 / 取消 ──
        bot = tk.Frame(self.dlg, bg=BG, pady=8)
        bot.pack(fill=tk.X, padx=12)

        self.status_var = tk.StringVar(value="")
        tk.Label(bot, textvariable=self.status_var, bg=BG, fg="#9cdcfe",
                 font=("Segoe UI", 9)).pack(side=tk.LEFT)

        tk.Button(bot, text="取消", command=self.dlg.destroy,
                  bg=BTN_BG, fg=FG, relief=tk.FLAT, font=FONT,
                  padx=12, cursor="hand2").pack(side=tk.RIGHT)
        self.exec_btn = tk.Button(bot, text="執行重命名", command=self._execute,
                                   bg="#4a7c59", fg="white", relief=tk.FLAT,
                                   font=("Segoe UI", 10, "bold"),
                                   padx=14, cursor="hand2")
        self.exec_btn.pack(side=tk.RIGHT, padx=6)

    def _on_wheel(self, event):
        self.lc.yview_scroll(-1 if event.delta > 0 else 1, "units")

    def _select_all(self):
        for v in self.check_vars:
            v.set(True)
        self._update_preview()

    def _deselect_all(self):
        for v in self.check_vars:
            v.set(False)
        self._update_preview()

    def _update_preview(self):
        prefix = self.prefix_var.get().strip()
        try:
            n = int(self.start_var.get())
        except ValueError:
            n = 1

        selected_count = sum(v.get() for v in self.check_vars)
        self.count_var.set(f"已選 {selected_count} / {len(self.files)} 個")

        for i, fpath in enumerate(self.files):
            if self.check_vars[i].get():
                ext = os.path.splitext(fpath)[1]
                self.preview_vars[i].set(f"→  {prefix}_{n}{ext}")
                n += 1
            else:
                self.preview_vars[i].set("(跳過)")

    def _execute(self):
        prefix = self.prefix_var.get().strip()
        if not prefix:
            messagebox.showwarning("警告", "請輸入前綴名稱", parent=self.dlg)
            return

        # 安全性：移除路徑分隔符，防止路徑遍歷攻擊
        prefix = os.path.basename(prefix)
        invalid_chars = r'\/:*?"<>|'
        if any(c in prefix for c in invalid_chars):
            messagebox.showwarning("警告", f"前綴不可包含以下字元：{invalid_chars}", parent=self.dlg)
            return
        if not prefix:
            messagebox.showwarning("警告", "請輸入有效的前綴名稱", parent=self.dlg)
            return
        try:
            start = int(self.start_var.get())
        except ValueError:
            messagebox.showwarning("警告", "起始數字必須是整數", parent=self.dlg)
            return

        targets = [(i, self.files[i]) for i, v in enumerate(self.check_vars) if v.get()]
        if not targets:
            messagebox.showwarning("警告", "請至少選擇一個檔案", parent=self.dlg)
            return

        # 計算最終目標路徑
        plan = []   # (old_path, final_path, final_name)
        n = start
        for _, old_path in targets:
            ext        = os.path.splitext(old_path)[1]
            final_name = f"{prefix}_{n}{ext}"
            # 保留原始目錄，支援遞迴掃描的子目錄檔案
            final_path = os.path.join(os.path.dirname(old_path), final_name)
            plan.append((old_path, final_path, final_name))
            n += 1

        # 衝突檢查：目標存在且不屬於本次重命名來源
        old_set = {p for p, _, _ in plan}
        conflicts = [fn for _, fp, fn in plan
                     if os.path.exists(fp) and fp not in old_set]
        if conflicts:
            preview = "\n".join(conflicts[:6])
            if len(conflicts) > 6:
                preview += f"\n...共 {len(conflicts)} 個"
            if not messagebox.askyesno("確認覆蓋",
                                        f"以下檔案將被覆蓋:\n{preview}\n\n確定繼續?",
                                        parent=self.dlg):
                return

        # 兩段式重命名：先 → 暫存名，再 → 最終名（避免中途衝突）
        tmp_plan = []
        errors = []
        for old_path, final_path, final_name in plan:
            ext  = os.path.splitext(old_path)[1]
            # Create temporary file in the same directory as the source file to
            # ensure moves are on the same filesystem and reduce cross-disk failures.
            tmp_dir = os.path.dirname(old_path) or self.folder_path
            tmp  = os.path.join(tmp_dir, f"_tmp_{uuid.uuid4().hex}{ext}")
            try:
                # use os.replace to be able to overwrite existing tmp if necessary
                os.replace(old_path, tmp)
                tmp_plan.append((tmp, final_path, final_name))
            except Exception as e:
                errors.append(f"{os.path.basename(old_path)}: {e}")

        renamed = 0
        for tmp, final_path, final_name in tmp_plan:
            try:
                # final move: replace if target exists
                os.replace(tmp, final_path)
                renamed += 1
            except Exception as e:
                errors.append(f"(暫存) → {final_name}: {e}")

        if errors:
            messagebox.showerror("部分失敗",
                                  f"完成 {renamed} 個，失敗:\n" + "\n".join(errors),
                                  parent=self.dlg)
        else:
            self.status_var.set(f"✓ 已成功重命名 {renamed} 個檔案")
            self.exec_btn.config(state=tk.DISABLED)

        self.on_done()


if __name__ == "__main__":
    root = TkinterDnD.Tk() if DND_AVAILABLE else tk.Tk()
    app = FileRenamer(root)
    root.mainloop()